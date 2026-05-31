import io
import base64
import re
from html import unescape
from html.parser import HTMLParser
from pathlib import Path
from typing import Any, Dict, Optional
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement

class DesignSystem:
    # Default design constants inspired by the reference script
    DEFAULT_FONT = "Arial"
    ACCENT_COLOR = RGBColor(0, 51, 102)  # Dark Blue
    TEXT_COLOR = RGBColor(44, 62, 80)
    BG_COLOR = RGBColor(255, 255, 255)


class HtmlTextRunParser(HTMLParser):
    BLOCK_TAGS = {"div", "p", "li", "section", "article", "blockquote", "h1", "h2", "h3", "h4", "h5", "h6"}

    def __init__(self):
        super().__init__(convert_charrefs=True)
        self.runs: list[dict[str, Any]] = []
        self._style_stack: list[dict[str, Any]] = [{"bold": False, "italic": False, "underline": False}]
        self._tag_stack: list[tuple[str, bool]] = []

    @property
    def current_style(self) -> dict[str, Any]:
        return self._style_stack[-1]

    def _push(self, **updates):
        self._style_stack.append({**self.current_style, **updates})

    def _mark_tag(self, tag: str, pushed: bool):
        self._tag_stack.append((tag, pushed))

    def _pop_tag(self, tag: str) -> bool:
        if not self._tag_stack:
            return False
        if self._tag_stack[-1][0] == tag:
            _, pushed = self._tag_stack.pop()
        else:
            pushed = False
            for idx in range(len(self._tag_stack) - 1, -1, -1):
                if self._tag_stack[idx][0] == tag:
                    _, pushed = self._tag_stack.pop(idx)
                    break
        if pushed and len(self._style_stack) > 1:
            self._style_stack.pop()
        return pushed

    def _newline(self):
        if self.runs and self.runs[-1].get("text") == "\n":
            return
        self.runs.append({"text": "\n", **self.current_style})

    def _parse_inline_style(self, style_text: str) -> dict[str, Any]:
        updates: dict[str, Any] = {}
        for declaration in str(style_text or "").split(";"):
            if ":" not in declaration:
                continue
            prop, value = declaration.split(":", 1)
            prop = prop.strip().lower()
            value = value.strip()
            if not value:
                continue
            if prop == "font-size":
                updates["fontSize"] = value
            elif prop == "font-family":
                updates["fontFamily"] = value
            elif prop == "color":
                updates["color"] = value
            elif prop == "font-weight":
                updates["fontWeight"] = value
                if value.lower() == "bold" or re.search(r"\d+", value) and int(re.search(r"\d+", value).group(0)) >= 600:
                    updates["bold"] = True
            elif prop == "font-style":
                updates["fontStyle"] = value
                if value.lower() == "italic":
                    updates["italic"] = True
            elif prop == "text-decoration" and "underline" in value.lower():
                updates["underline"] = True
            elif prop == "vertical-align":
                if value.lower() in {"super", "sup", "sub"}:
                    updates["verticalAlign"] = value.lower()
        return updates

    def _attrs_to_updates(self, attrs) -> dict[str, Any]:
        attrs_map = {str(key).lower(): value for key, value in attrs}
        updates = self._parse_inline_style(attrs_map.get("style", ""))
        if attrs_map.get("face"):
            updates["fontFamily"] = attrs_map["face"]
        if attrs_map.get("color"):
            updates["color"] = attrs_map["color"]
        if attrs_map.get("size"):
            # Legacy <font size="1..7"> values are approximate browser defaults.
            size_map = {"1": "10px", "2": "13px", "3": "16px", "4": "18px", "5": "24px", "6": "32px", "7": "48px"}
            updates["fontSize"] = size_map.get(str(attrs_map["size"]).strip(), attrs_map["size"])
        return updates

    def handle_starttag(self, tag, attrs):
        tag = tag.lower()
        attr_updates = self._attrs_to_updates(attrs)
        if tag in {"strong", "b"}:
            self._push(**attr_updates, bold=True)
            self._mark_tag(tag, True)
        elif tag in {"em", "i"}:
            self._push(**attr_updates, italic=True)
            self._mark_tag(tag, True)
        elif tag == "u":
            self._push(**attr_updates, underline=True)
            self._mark_tag(tag, True)
        elif tag == "sup":
            self._push(**attr_updates, verticalAlign="super")
            self._mark_tag(tag, True)
        elif tag == "sub":
            self._push(**attr_updates, verticalAlign="sub")
            self._mark_tag(tag, True)
        elif tag == "br":
            self._newline()
        elif tag in {"span", "font"}:
            self._push(**attr_updates)
            self._mark_tag(tag, True)
        elif tag in self.BLOCK_TAGS:
            if self.runs:
                self._newline()
            if attr_updates:
                self._push(**attr_updates)
                self._mark_tag(tag, True)
            else:
                self._mark_tag(tag, False)
        elif attr_updates:
            self._push(**attr_updates)
            self._mark_tag(tag, True)
        else:
            self._mark_tag(tag, False)

    def handle_endtag(self, tag):
        tag = tag.lower()
        if tag in {"strong", "b", "em", "i", "u", "span", "font", "sup", "sub"}:
            self._pop_tag(tag)
        elif tag in self.BLOCK_TAGS:
            self._pop_tag(tag)
            self._newline()
        else:
            self._pop_tag(tag)

    def handle_data(self, data):
        if data:
            self.runs.append({"text": unescape(data), **self.current_style})


def html_to_text_runs(content: Any) -> list[dict[str, Any]]:
    parser = HtmlTextRunParser()
    parser.feed(str(content or ""))
    runs = []
    for run in parser.runs:
        text = run.get("text", "")
        if text == "\n":
            if runs and runs[-1].get("text") == "\n":
                continue
            runs.append(run)
        elif text:
            runs.append(run)
    while runs and runs[-1].get("text") == "\n":
        runs.pop()
    return runs or [{"text": ""}]

class PPTXExporter:
    def __init__(self, state: Dict[str, Any], project_root: Optional[Path | str] = None):
        self.state = state
        self.prs = Presentation()
        self.project_root = self._resolve_project_root(project_root)

        # Mapping from SlideForge page setup to PPTX dimensions
        setup = state.get("pageSetup", "standard-4-3")
        if setup == "widescreen-16-9":
            self.prs.slide_width = Inches(13.333)
            self.prs.slide_height = Inches(7.5)
            self.base_w, self.base_h = 1280, 720
        elif setup == "widescreen-16-10":
            self.prs.slide_width = Inches(12.8)
            self.prs.slide_height = Inches(8.0)
            self.base_w, self.base_h = 1280, 800
        else: # standard-4-3
            self.prs.slide_width = Inches(10.666)
            self.prs.slide_height = Inches(8.0)
            self.base_w, self.base_h = 1024, 768
        self.theme = self._theme_defaults(str(state.get("presentationTheme") or "editorial"))
        self.slide_bg_rgb = self.theme["background"]

    def _theme_defaults(self, theme_id: str) -> dict[str, RGBColor]:
        themes = {
            "editorial": ("#EEF2F7", "#2E2E2E", "#6B7280", "#3B82F6", "#A7C7E7"),
            "blueprint": ("#F1F7FD", "#10233B", "#5F7287", "#2563EB", "#0F766E"),
            "fieldnotes": ("#EFE4D0", "#2F261D", "#6F6559", "#7A8F47", "#B48A4A"),
            "monograph": ("#ECEFF3", "#171717", "#6B7280", "#1F2937", "#9CA3AF"),
            "graphite": ("#111827", "#F5F7FB", "#94A3B8", "#22D3EE", "#38BDF8"),
            "horizon": ("#0B2442", "#EEF4FF", "#AFBDD6", "#7DD3FC", "#4F7CFF"),
            "chalkboard": ("#10261F", "#F8F3E7", "#D4CBB7", "#F5D76E", "#8ED1C7"),
            "circuit": ("#0A2320", "#ECF7F5", "#9AB8B3", "#63E6D8", "#1FB6A6"),
            "afterglow": ("#1A2034", "#F5F7FF", "#BDC4DE", "#F3B76A", "#6E86FF"),
            "sage": ("#EDF4E8", "#25342D", "#68786E", "#4F7D5A", "#B7A66C"),
            "porcelain": ("#EEF7F9", "#233047", "#6C7890", "#477CA4", "#C08393"),
            "rosewater": ("#F7E7E4", "#3A2B32", "#7F6B72", "#B76E79", "#AA8650"),
            "buttercup": ("#F5EBC5", "#332B15", "#756C4C", "#D59D20", "#5F9F76"),
            "tidepool": ("#E4F4F3", "#17323A", "#5F7880", "#168B8B", "#4B8FC5"),
            "lavender": ("#EFEAF8", "#302C45", "#756F8D", "#7D6BB3", "#CF8A72"),
            "midnightGarden": ("#10231A", "#EFF7ED", "#B7C9B7", "#9AE6B4", "#5FAF79"),
            "retroPop": ("#FFE7DF", "#28212A", "#745F76", "#EF476F", "#06D6A0"),
        }
        bg, text, muted, accent, accent2 = themes.get(theme_id, themes["editorial"])
        return {
            "background": self._parse_color(bg, DesignSystem.BG_COLOR),
            "text": self._parse_color(text, DesignSystem.TEXT_COLOR),
            "muted": self._parse_color(muted, RGBColor(100, 116, 139)),
            "accent": self._parse_color(accent, DesignSystem.ACCENT_COLOR),
            "accent2": self._parse_color(accent2, DesignSystem.ACCENT_COLOR),
        }

    def _resolve_project_root(self, project_root: Optional[Path | str] = None) -> Path:
        if project_root:
            return Path(project_root).expanduser().resolve()
        try:
            from django.conf import settings

            base_dir = getattr(settings, "BASE_DIR", None)
            if base_dir:
                return Path(base_dir).expanduser().resolve()
        except Exception:
            pass
        return Path(__file__).resolve().parent.parent

    def _safe_asset_roots(self) -> list[Path]:
        roots = [
            self.project_root / "media",
            self.project_root / "static",
            self.project_root / "staticfiles",
            self.project_root / "assets",
        ]
        try:
            from django.conf import settings

            for setting_name in ("MEDIA_ROOT", "STATIC_ROOT"):
                configured = getattr(settings, setting_name, None)
                if configured:
                    roots.append(Path(configured))
            for static_dir in getattr(settings, "STATICFILES_DIRS", []) or []:
                roots.append(Path(static_dir))
        except Exception:
            pass

        resolved = []
        seen = set()
        for root in roots:
            try:
                path = Path(root).expanduser().resolve()
            except Exception:
                continue
            if path in seen:
                continue
            seen.add(path)
            resolved.append(path)
        return resolved

    def _resolve_safe_local_image_path(self, content: str) -> Optional[Path]:
        normalized = str(content or "").strip()
        if not normalized or "://" in normalized:
            return None

        candidates: list[Path] = []
        if normalized.startswith("/media/"):
            candidates.append(self.project_root / normalized.lstrip("/"))
        elif normalized.startswith("/static/"):
            candidates.append(self.project_root / normalized.lstrip("/"))
            candidates.append(self.project_root / "staticfiles" / normalized.removeprefix("/static/"))
        elif normalized.startswith("/assets/"):
            candidates.append(self.project_root / normalized.lstrip("/"))
        elif normalized.startswith("/"):
            candidates.append(Path(normalized))
        else:
            candidates.append(self.project_root / normalized)

        safe_roots = self._safe_asset_roots()
        for candidate in candidates:
            try:
                path = candidate.expanduser().resolve()
            except Exception:
                continue
            if not path.exists() or not path.is_file():
                continue
            if any(path == root or path.is_relative_to(root) for root in safe_roots):
                return path
        return None

    def _parse_font_weight(self, value: Any) -> int:
        if isinstance(value, (int, float)):
            return int(value)
        normalized = str(value or "400").strip().lower()
        if normalized == "normal":
            return 400
        if normalized == "bold":
            return 700
        match = re.search(r"\d+", normalized)
        if match:
            return int(match.group(0))
        return 400

    def _px_to_inches_w(self, px: float) -> Inches:
        # Scale pixels relative to the logical canvas width
        return Inches((px / self.base_w) * (self.prs.slide_width / Inches(1)))

    def _px_to_inches_h(self, px: float) -> Inches:
        # Scale pixels relative to the logical canvas height
        return Inches((px / self.base_h) * (self.prs.slide_height / Inches(1)))

    def _parse_color_components(self, color_str: Optional[str]) -> Optional[tuple[int, int, int, float]]:
        if not color_str:
            return None
        normalized = str(color_str).strip()
        if not normalized or normalized.lower() in {"transparent", "none", "inherit", "currentcolor"}:
            return None
        try:
            if normalized.startswith("#"):
                h = normalized.lstrip("#")
                if len(h) == 3:
                    h = "".join([c*2 for c in h])
                if len(h) >= 6:
                    return int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16), 1.0
            elif normalized.startswith("rgb"):
                matches = re.findall(r"[\d.]+", normalized)
                if len(matches) >= 3:
                    alpha = float(matches[3]) if len(matches) >= 4 else 1.0
                    return (
                        max(0, min(255, int(float(matches[0])))),
                        max(0, min(255, int(float(matches[1])))),
                        max(0, min(255, int(float(matches[2])))),
                        max(0.0, min(1.0, alpha)),
                    )
        except Exception:
            return None
        return None

    def _composite_color(self, color: tuple[int, int, int, float], bg: Optional[RGBColor] = None) -> RGBColor:
        r, g, b, alpha = color
        if alpha >= 0.995:
            return RGBColor(r, g, b)
        bg = bg or self.slide_bg_rgb
        return RGBColor(
            round(r * alpha + bg[0] * (1 - alpha)),
            round(g * alpha + bg[1] * (1 - alpha)),
            round(b * alpha + bg[2] * (1 - alpha)),
        )

    def _parse_color(self, color_str: Optional[str], default: RGBColor = DesignSystem.TEXT_COLOR) -> RGBColor:
        parsed = self._parse_color_components(color_str)
        if parsed:
            return self._composite_color(parsed)
        return default

    def _is_transparent(self, color_str: Optional[str]) -> bool:
        normalized = str(color_str or "").strip().lower()
        if normalized in {"", "transparent", "none"}:
            return True
        parsed = self._parse_color_components(normalized)
        return bool(parsed and parsed[3] <= 0.01)

    def _parse_px(self, value: Any, fallback: float = 100) -> float:
        match = re.search(r"-?\d+(?:\.\d+)?", str(value or ""))
        if not match:
            return fallback
        try:
            return float(match.group(0))
        except ValueError:
            return fallback

    def _set_text_common(self, run: Any, styles: dict[str, Any], *, extra: Optional[dict[str, bool]] = None):
        extra = extra or {}
        font_size = int(self._parse_px(extra.get("fontSize") or styles.get("fontSize", "18"), 18))
        font_color = self._parse_color(extra.get("color") or styles.get("color"), self.theme["text"])
        font_name = str(extra.get("fontFamily") or styles.get("fontFamily", DesignSystem.DEFAULT_FONT)).split(",")[0].replace('"', "").strip()
        run.font.size = Pt(font_size)
        run.font.color.rgb = font_color
        run.font.name = font_name
        run.font.bold = self._parse_font_weight(extra.get("fontWeight") or styles.get("fontWeight")) >= 600 or extra.get("bold", False)
        run.font.italic = (extra.get("fontStyle") or styles.get("fontStyle")) == "italic" or extra.get("italic", False)
        run.font.underline = bool(extra.get("underline", False))
        if extra.get("verticalAlign") == "super" and hasattr(run.font, "superscript"):
            run.font.superscript = True
        if extra.get("verticalAlign") == "sub" and hasattr(run.font, "subscript"):
            run.font.subscript = True

    def _paragraph_alignment(self, styles: dict[str, Any]):
        align = str(styles.get("textAlign", "left")).lower()
        if align == "center":
            return PP_ALIGN.CENTER
        if align == "right":
            return PP_ALIGN.RIGHT
        if align == "justify":
            return PP_ALIGN.JUSTIFY
        return PP_ALIGN.LEFT

    def _get_image_stream(self, content: str) -> Optional[io.BytesIO]:
        if not content:
            return None

        if content.startswith("data:"):
            try:
                base64_data = content.split(",")[1]
                return io.BytesIO(base64.b64decode(base64_data))
            except Exception:
                return None

        try:
            path = self._resolve_safe_local_image_path(content)
            if path:
                with open(path, "rb") as f:
                    return io.BytesIO(f.read())
        except Exception as e:
            print(f"Error resolving image path {content}: {e}")

        return None

    def _apply_fill(self, shape: Any, color: Optional[str], default: Optional[RGBColor] = None):
        if self._is_transparent(color):
            shape.fill.background()
            return
        shape.fill.solid()
        shape.fill.fore_color.rgb = self._parse_color(color, default or self.theme["accent"])

    def _apply_line(self, shape: Any, styles: dict[str, Any], default_color: Optional[RGBColor] = None):
        border = str(styles.get("border") or "").strip()
        color = styles.get("borderColor") or styles.get("stroke")
        width = styles.get("borderWidth") or styles.get("strokeWidth")
        if border:
            color_match = re.search(r"#[0-9a-fA-F]{3,8}|rgba?\([^)]+\)", border)
            if color_match:
                color = color_match.group(0)
            width_match = re.search(r"(-?\d+(?:\.\d+)?)px", border)
            if width_match:
                width = width_match.group(1)
        if not color and not width:
            shape.line.fill.background()
            return
        if self._is_transparent(color):
            shape.line.fill.background()
            return
        shape.line.color.rgb = self._parse_color(color, default_color or self.theme["accent"])
        shape.line.width = Pt(max(0.25, float(self._parse_px(width, 1))))

    def _get_bullet_char(self, style: str, level: int) -> str:
        """Get the appropriate bullet character for the given style and level."""
        # CRITICAL FIX #5: Bullet character mapping based on SlideForge theme
        bullet_chars = {
            "default": ["•", "◦", "▪"],
            "square": ["■", "□", "▪"],
            "diamond": ["◆", "◇", "◈"],
            "modern": ["→", "–", "◆"],
            "chevron": ["»", "›", "–"],
            "dash": ["–", "—", "·"],
            "checklist": ["✓", "✓", "✓"],
            "star": ["✦", "✧", "•"],
        }
        chars = bullet_chars.get(style, bullet_chars["default"])
        return chars[min(level, len(chars) - 1)]

    def _set_paragraph_bullet(self, paragraph: Any, char: str = "•"):
        try:
            p_pr = paragraph._p.get_or_add_pPr()
            for child in list(p_pr):
                if child.tag.endswith("}buNone") or child.tag.endswith("}buChar") or child.tag.endswith("}buAutoNum"):
                    p_pr.remove(child)

            # Use native PowerPoint auto-numbering if char is special indicator
            if char == "__AUTO_BULLET__":
                # Let PowerPoint handle bullets using default formatting
                bullet = OxmlElement("a:buFont")
                bullet.set("typeface", "Arial")
                p_pr.append(bullet)
            elif char == "__AUTO_NUMBER__":
                # Use native PowerPoint numbering
                buAutoNum = OxmlElement("a:buAutoNum")
                buAutoNum.set("type", "alphaLcParenBoth")
                p_pr.append(buAutoNum)
            else:
                # Custom bullet character
                bullet = OxmlElement("a:buChar")
                bullet.set("char", char)
                p_pr.append(bullet)
        except Exception:
            return

    def export(self) -> io.BytesIO:
        slides = self.state.get("slides", [])
        for slide_data in slides:
            self._add_slide(slide_data)

        output = io.BytesIO()
        self.prs.save(output)
        output.seek(0)
        return output

    def _add_slide(self, slide_data: Dict[str, Any]):
        self.slide_bg_rgb = self.theme["background"]
        # Use a blank layout (index 6 is usually blank)
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        # 1. Background
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = self.slide_bg_rgb
        bg = slide_data.get("background")
        if bg and bg.get("content"):
            if bg.get("type") == "color":
                fill.fore_color.rgb = self._parse_color(bg["content"], self.slide_bg_rgb)
                self.slide_bg_rgb = fill.fore_color.rgb
            elif bg.get("type") == "image":
                img_stream = self._get_image_stream(bg["content"])
                if img_stream:
                    slide.shapes.add_picture(img_stream, 0, 0, width=self.prs.slide_width, height=self.prs.slide_height)
        elif slide_data.get("backgroundColor"):
            fill.fore_color.rgb = self._parse_color(slide_data.get("backgroundColor"), self.slide_bg_rgb)
            self.slide_bg_rgb = fill.fore_color.rgb

        if slide_data.get("notes"):
            slide.notes_slide.notes_text_frame.text = str(slide_data.get("notes") or "")

        # 2. Elements
        elements = [*self._build_master_elements(slide_data), *slide_data.get("elements", [])]
        # Sort by zIndex
        sorted_elements = sorted(elements, key=lambda e: e.get("styles", {}).get("zIndex", 0))

        for el in sorted_elements:
            self._add_element(slide, el)

    def _master_config(self, master_id: str) -> Optional[dict[str, Any]]:
        if master_id == "none":
            return None
        defaults = {
            "content": {
                "id": "content",
                "enabled": True,
                "footerText": "Presentation",
                "logoText": "SlideForge",
                "showSlideNumber": True,
                "showFooter": True,
                "showTopRule": True,
            },
            "title": {
                "id": "title",
                "enabled": True,
                "footerText": "",
                "logoText": "SlideForge",
                "showSlideNumber": False,
                "showFooter": False,
                "showTopRule": True,
            },
            "section": {
                "id": "section",
                "enabled": True,
                "footerText": "Section",
                "logoText": "SlideForge",
                "showSlideNumber": True,
                "showFooter": True,
                "showTopRule": False,
            },
        }
        master_id = master_id if master_id in defaults else "content"
        masters = self.state.get("masterSlides") if isinstance(self.state.get("masterSlides"), dict) else {}
        config = {**defaults[master_id], **(masters.get(master_id) if isinstance(masters.get(master_id), dict) else {})}
        return config if config.get("enabled", True) else None

    def _master_shape(self, x: float, y: float, w: float, h: float, color: str, *, opacity: Optional[str] = None, border: Optional[str] = None, radius: str = "0px") -> dict[str, Any]:
        return {
            "type": "shape",
            "shapeType": "rectangle",
            "x": x,
            "y": y,
            "width": f"{w}px",
            "height": f"{h}px",
            "styles": {
                "backgroundColor": color,
                "borderRadius": radius,
                "zIndex": -20,
                **({"opacity": opacity} if opacity is not None else {}),
                **({"border": border} if border else {}),
            },
        }

    def _master_text(self, x: float, y: float, w: float, text: str, styles: dict[str, Any]) -> dict[str, Any]:
        return {
            "type": "text",
            "x": x,
            "y": y,
            "width": f"{w}px",
            "height": "28px",
            "content": text,
            "styles": {"zIndex": -19, "backgroundColor": "transparent", **styles},
        }

    def _build_master_elements(self, slide_data: dict[str, Any]) -> list[dict[str, Any]]:
        master_id = str(slide_data.get("masterId") or "content")
        config = self._master_config(master_id)
        if not config:
            return []
        try:
            slide_index = self.state.get("slides", []).index(slide_data)
        except ValueError:
            slide_index = 0
        slide_number = str(slide_index + 1).zfill(2)
        sx = self.base_w / 1024
        sy = self.base_h / 768
        text = self._rgb_to_hex(self.theme["text"])
        muted = self._rgb_to_hex(self.theme["muted"])
        accent = self._rgb_to_hex(self.theme["accent"])
        accent2 = self._rgb_to_hex(self.theme["accent2"])
        border = self._rgb_to_rgba(self.theme["muted"], 0.28)
        surface = self._rgb_to_rgba(self.theme["text"], 0.06)
        heading_font = "Aptos Display"
        body_font = DesignSystem.DEFAULT_FONT

        def sc(el: dict[str, Any]) -> dict[str, Any]:
            next_el = {**el, "x": round(float(el.get("x", 0)) * sx), "y": round(float(el.get("y", 0)) * sy)}
            next_el["width"] = f"{max(1, self._parse_px(el.get('width'), 1) * sx)}px"
            next_el["height"] = f"{max(1, self._parse_px(el.get('height'), 1) * sy)}px"
            styles = {**next_el.get("styles", {})}
            if styles.get("fontSize"):
                styles["fontSize"] = f"{max(1, self._parse_px(styles.get('fontSize'), 1) * min(sx, sy))}px"
            next_el["styles"] = styles
            return next_el

        elements: list[dict[str, Any]] = []
        if master_id == "title":
            if config.get("showTopRule", True):
                elements.append(self._master_shape(0, 0, 1024, 8, accent))
            elements.append(self._master_shape(882, 48, 88, 88, accent2, opacity="0.18", radius="22px"))
            elements.append(self._master_shape(930, 96, 40, 40, accent, opacity="0.82", radius="999px"))
            if config.get("logoText"):
                elements.append(self._master_text(64, 702, 420, str(config.get("logoText")), {
                    "color": muted, "fontFamily": body_font, "fontSize": "12px", "fontWeight": "700",
                }))
        elif master_id == "section":
            elements.append(self._master_shape(0, 0, 10, 768, accent))
            elements.append(self._master_shape(24, 42, 92, 92, surface, border=f"1px solid {border}", radius="24px"))
            elements.append(self._master_text(47, 67, 60, slide_number, {
                "color": accent, "fontFamily": heading_font, "fontSize": "34px", "fontWeight": "800", "textAlign": "center",
            }))
        else:
            if config.get("showTopRule", True):
                elements.append(self._master_shape(0, 0, 1024, 5, accent))
            elements.append(self._master_shape(52, 712, 920, 1, border))

        if config.get("showFooter", True):
            if config.get("logoText"):
                elements.append(self._master_text(54, 724, 145, str(config.get("logoText")), {
                    "color": text, "fontFamily": body_font, "fontSize": "11px", "fontWeight": "800",
                }))
            if config.get("footerText"):
                elements.append(self._master_text(224, 724, 520, str(config.get("footerText")), {
                    "color": muted, "fontFamily": body_font, "fontSize": "11px", "fontWeight": "600",
                }))
        if config.get("showSlideNumber", True):
            elements.append(self._master_text(918, 724, 54, slide_number, {
                "color": muted, "fontFamily": body_font, "fontSize": "11px", "fontWeight": "800", "textAlign": "right",
            }))
        return [sc(el) for el in elements]

    def _rgb_to_hex(self, color: RGBColor) -> str:
        return f"#{color[0]:02X}{color[1]:02X}{color[2]:02X}"

    def _rgb_to_rgba(self, color: RGBColor, alpha: float) -> str:
        return f"rgba({color[0]},{color[1]},{color[2]},{alpha})"

    def _add_element(self, slide: Any, el: Dict[str, Any]):
        try:
            x = self._px_to_inches_w(float(el.get("x", 0)))
            y = self._px_to_inches_h(float(el.get("y", 0)))

            # Width and height can be "auto" or strings like "400px"
            w_raw = el.get("width", "100")
            h_raw = el.get("height", "100")

            w = self._px_to_inches_w(float(re.sub(r"[^\d.]", "", str(w_raw)) or 100))
            h = self._px_to_inches_h(float(re.sub(r"[^\d.]", "", str(h_raw)) or 100))

            el_type = el.get("type")
            if el_type == "text":
                self._add_text_element(slide, el, x, y, w, h)
            elif el_type == "image":
                self._add_image_element(slide, el, x, y, w, h)
            elif el_type == "shape":
                self._add_shape_element(slide, el, x, y, w, h)
            elif el_type == "table":
                self._add_table_element(slide, el, x, y, w, h)
            elif el_type == "connector":
                self._add_connector_element(slide, el, x, y, w, h)
            elif el_type == "video":
                self._add_video_element(slide, el, x, y, w, h)
            elif el_type == "mermaid":
                self._add_mermaid_element(slide, el, x, y, w, h)
            elif el_type in {"html", "pdf", "molecule", "chart", "equation", "latex"}:
                self._add_placeholder_element(slide, el, x, y, w, h)
        except Exception as e:
            print(f"Error adding element {el.get('id')}: {e}")

    def _add_text_element(self, slide: Any, el: Dict[str, Any], x: Inches, y: Inches, w: Inches, h: Inches):
        styles = el.get("styles", {})
        content = el.get("content", "")

        # Create textbox
        shape = slide.shapes.add_textbox(x, y, w, h)
        if styles.get("backgroundColor") and not self._is_transparent(styles.get("backgroundColor")):
            self._apply_fill(shape, styles.get("backgroundColor"), self.slide_bg_rgb)
        else:
            shape.fill.background()
        self._apply_line(shape, styles, self.theme["accent"])
        tf = shape.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        # Handle structured bullet points or raw HTML
        tf.clear()
        if isinstance(content, list):
            # CRITICAL FIX #5: Enhanced list support for PPTX export
            for idx, item in enumerate(content):
                p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                p.level = max(0, min(8, int(item.get("level", 0) or 0)))
                p.alignment = self._paragraph_alignment(styles)
                
                # Set bullet based on list kind and style
                kind = item.get("kind", "bullet")
                style = item.get("style", "default")
                
                if kind == "numbered":
                    # Use auto-numbered format for ordered lists
                    self._set_paragraph_bullet(p, "__AUTO_NUMBER__")
                else:
                    # Use bullet character based on style
                    char = self._get_bullet_char(style, p.level)
                    self._set_paragraph_bullet(p, char)
                
                raw = item.get("html") if isinstance(item, dict) and item.get("html") is not None else item.get("text", "")
                self._append_runs_to_paragraph(p, html_to_text_runs(raw), styles)
        else:
            p = tf.paragraphs[0]
            p.alignment = self._paragraph_alignment(styles)
            self._append_runs_to_paragraph(p, html_to_text_runs(content), styles)

    def _append_runs_to_paragraph(self, paragraph: Any, runs: list[dict[str, Any]], styles: dict[str, Any]):
        current = paragraph
        for run_data in runs:
            text = run_data.get("text", "")
            parts = text.split("\n")
            for idx, part in enumerate(parts):
                if idx > 0:
                    current = paragraph._parent.add_paragraph()
                    current.alignment = paragraph.alignment
                if not part:
                    continue
                run = current.add_run()
                run.text = part
                self._set_text_common(run, styles, extra=run_data)

    def _add_image_element(self, slide: Any, el: Dict[str, Any], x: Inches, y: Inches, w: Inches, h: Inches):
        content = el.get("content", "")
        img_stream = self._get_image_stream(content)
        if not img_stream:
            return

        pic = slide.shapes.add_picture(img_stream, x, y, width=w, height=h)
        self._apply_image_crop_transform(pic, el.get("cropTransform"))

    def _apply_image_crop_transform(self, pic: Any, crop: Optional[dict[str, Any]]):
        if not crop:
            return
        try:
            width_pct = max(100.0, float(crop.get("widthPercent", 100) or 100))
            height_pct = max(100.0, float(crop.get("heightPercent", 100) or 100))
            left_pct = min(0.0, max(100.0 - width_pct, float(crop.get("leftPercent", 0) or 0)))
            top_pct = min(0.0, max(100.0 - height_pct, float(crop.get("topPercent", 0) or 0)))
        except (TypeError, ValueError):
            return

        pic.crop_left = max(0.0, min(1.0, -left_pct / width_pct))
        pic.crop_right = max(0.0, min(1.0, (width_pct + left_pct - 100.0) / width_pct))
        pic.crop_top = max(0.0, min(1.0, -top_pct / height_pct))
        pic.crop_bottom = max(0.0, min(1.0, (height_pct + top_pct - 100.0) / height_pct))

    def _add_shape_element(self, slide: Any, el: Dict[str, Any], x: Inches, y: Inches, w: Inches, h: Inches):
        shape_type_str = el.get("shapeType", "rectangle")

        mapping = {
            "rectangle": MSO_SHAPE.RECTANGLE,
            "circle": MSO_SHAPE.OVAL,
            "ellipse": MSO_SHAPE.OVAL,
            "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
            "diamond": MSO_SHAPE.DIAMOND,
            "hexagon": MSO_SHAPE.HEXAGON,
            "parallelogram": MSO_SHAPE.PARALLELOGRAM,
            "arrow-right": MSO_SHAPE.RIGHT_ARROW,
            "arrow-left": MSO_SHAPE.LEFT_ARROW,
            "arrow-up": MSO_SHAPE.UP_ARROW,
            "arrow-down": MSO_SHAPE.DOWN_ARROW,
        }

        mso_type = mapping.get(shape_type_str, MSO_SHAPE.RECTANGLE)
        shape = slide.shapes.add_shape(mso_type, x, y, w, h)

        styles = el.get("styles", {})
        self._apply_fill(shape, styles.get("backgroundColor"), self.theme["accent"])
        self._apply_line(shape, styles, self.theme["accent"])

    def _add_table_element(self, slide: Any, el: Dict[str, Any], x: Inches, y: Inches, w: Inches, h: Inches):
        table_data = el.get("tableData") or {}
        rows = max(1, int(table_data.get("rows") or len(table_data.get("cells") or []) or 1))
        cols = max(1, int(table_data.get("cols") or max((len(row) for row in table_data.get("cells", []) if isinstance(row, list)), default=1)))
        rows = min(rows, 50)
        cols = min(cols, 20)
        table_shape = slide.shapes.add_table(rows, cols, x, y, w, h)
        table = table_shape.table
        cells = table_data.get("cells") if isinstance(table_data.get("cells"), list) else []
        header_row = table_data.get("headerRow", True)
        font_size = table_data.get("fontSize") or "14px"

        for col_idx, raw_width in enumerate(table_data.get("colWidths") or []):
            if col_idx < cols:
                table.columns[col_idx].width = self._px_to_inches_w(max(36, float(raw_width or 140)))
        for row_idx, raw_height in enumerate(table_data.get("rowHeights") or []):
            if row_idx < rows:
                table.rows[row_idx].height = self._px_to_inches_h(max(24, float(raw_height or 44)))

        for row_idx in range(rows):
            for col_idx in range(cols):
                cell = table.cell(row_idx, col_idx)
                raw_cell = cells[row_idx][col_idx] if row_idx < len(cells) and isinstance(cells[row_idx], list) and col_idx < len(cells[row_idx]) else {}
                text = raw_cell.get("text", "") if isinstance(raw_cell, dict) else str(raw_cell or "")
                cell_styles = raw_cell.get("styles", {}) if isinstance(raw_cell, dict) and isinstance(raw_cell.get("styles"), dict) else {}
                cell.text = str(text)
                if row_idx == 0 and header_row:
                    fill_color = cell_styles.get("backgroundColor") or table_data.get("headerFill") or "#e2e8f0"
                elif table_data.get("zebra") and row_idx % 2 == 0:
                    fill_color = cell_styles.get("backgroundColor") or table_data.get("altFill") or "#f8fafc"
                else:
                    fill_color = cell_styles.get("backgroundColor") or table_data.get("bodyFill") or "#ffffff"
                cell.fill.solid()
                cell.fill.fore_color.rgb = self._parse_color(fill_color, self.slide_bg_rgb)
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.alignment = self._paragraph_alignment({"textAlign": cell_styles.get("textAlign") or table_data.get("textAlign", "left")})
                    for run in paragraph.runs:
                        self._set_text_common(
                            run,
                            {
                                "fontSize": cell_styles.get("fontSize") or font_size,
                                "fontFamily": cell_styles.get("fontFamily") or table_data.get("fontFamily", DesignSystem.DEFAULT_FONT),
                                "fontWeight": cell_styles.get("fontWeight") or ("700" if row_idx == 0 and header_row else table_data.get("fontWeight", "400")),
                                "fontStyle": cell_styles.get("fontStyle") or table_data.get("fontStyle"),
                                "color": cell_styles.get("color") or table_data.get("headerTextColor" if row_idx == 0 and header_row else "textColor") or ("#ffffff" if row_idx == 0 and header_row else None),
                            },
                        )
    def _add_connector_element(self, slide: Any, el: Dict[str, Any], x: Inches, y: Inches, w: Inches, h: Inches):
        styles = el.get("styles", {})
        connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x, y, x + w, y + h)
        connector.line.color.rgb = self._parse_color(styles.get("color") or styles.get("borderColor"), self.theme["accent"])
        connector.line.width = Pt(float(self._parse_px(styles.get("strokeWidth") or styles.get("borderWidth"), 2)))

    def _sanitize_svg(self, svg: str) -> str:
        cleaned = str(svg or "")
        cleaned = re.sub(r"<\s*(script|foreignObject)\b[^>]*>.*?<\s*/\s*\1\s*>", "", cleaned, flags=re.I | re.S)
        cleaned = re.sub(r"\s+on[a-zA-Z]+\s*=\s*(['\"]).*?\1", "", cleaned, flags=re.I | re.S)
        cleaned = re.sub(r"\s+(href|xlink:href)\s*=\s*(['\"])\s*(javascript:|data:text/html).*?\2", "", cleaned, flags=re.I | re.S)
        cleaned = re.sub(r"url\s*\(\s*(['\"]?)\s*(javascript:|data:text/html).*?\)", "none", cleaned, flags=re.I | re.S)
        return cleaned.strip()

    def _svg_to_png_stream(self, svg: str, width_px: int, height_px: int) -> Optional[io.BytesIO]:
        try:
            import cairosvg

            png = cairosvg.svg2png(
                bytestring=svg.encode("utf-8"),
                output_width=max(1, width_px),
                output_height=max(1, height_px),
            )
            stream = io.BytesIO(png)
            stream.seek(0)
            return stream
        except Exception:
            pass

        try:
            from svglib.svglib import svg2rlg
            from reportlab.graphics import renderPM

            drawing = svg2rlg(io.BytesIO(svg.encode("utf-8")))
            png = renderPM.drawToString(drawing, fmt="PNG", dpi=220)
            stream = io.BytesIO(png)
            stream.seek(0)
            return stream
        except Exception:
            return None

    def _add_mermaid_element(self, slide: Any, el: Dict[str, Any], x: Inches, y: Inches, w: Inches, h: Inches):
        svg = self._sanitize_svg(el.get("svgContent") or "")
        if not svg:
            self._add_placeholder_element(slide, {**el, "type": "mermaid"}, x, y, w, h)
            return

        svg_stream = io.BytesIO(svg.encode("utf-8"))
        try:
            slide.shapes.add_picture(svg_stream, x, y, width=w, height=h)
            return
        except Exception:
            pass

        width_px = int(max(320, self._parse_px(el.get("width"), 560) * 3))
        height_px = int(max(240, self._parse_px(el.get("height"), 360) * 3))
        png_stream = self._svg_to_png_stream(svg, width_px, height_px)
        if png_stream:
            slide.shapes.add_picture(png_stream, x, y, width=w, height=h)
            return
        self._add_placeholder_element(slide, {**el, "type": "mermaid"}, x, y, w, h)

    def _add_placeholder_element(self, slide: Any, el: Dict[str, Any], x: Inches, y: Inches, w: Inches, h: Inches):
        label_by_type = {
            "video": "Video",
            "html": "HTML embed",
            "pdf": "PDF",
            "molecule": "3D molecule",
            "chart": "Chart",
            "equation": "Equation",
            "latex": "Equation",
            "mermaid": "Mermaid diagram",
        }
        label = label_by_type.get(el.get("type"), "Unsupported content")
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = self._parse_color("rgba(148, 163, 184, 0.18)", self.theme["muted"])
        shape.line.color.rgb = self.theme["muted"]
        tf = shape.text_frame
        tf.text = f"{label} placeholder"
        paragraph = tf.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        if paragraph.runs:
            run = paragraph.runs[0]
        else:
            run = paragraph.add_run()
            run.text = f"{label} placeholder"
        run.font.size = Pt(14)
        run.font.color.rgb = self.theme["muted"]

    def _resolve_local_video_path(self, content: str) -> Optional[Path]:
        normalized = str(content or "").strip()
        if not normalized or "://" in normalized:
            return None

        candidates: list[Path] = []
        if normalized.startswith("/media/"):
            candidates.append(self.project_root / normalized.lstrip("/"))
        elif normalized.startswith("/static/"):
            candidates.append(self.project_root / normalized.lstrip("/"))
        elif normalized.startswith("/"):
            candidates.append(Path(normalized))
        else:
            candidates.append(self.project_root / normalized)

        safe_roots = self._safe_asset_roots()
        for candidate in candidates:
            try:
                path = candidate.expanduser().resolve()
            except Exception:
                continue
            if path.exists() and path.is_file():
                if any(path == root or path.is_relative_to(root) for root in safe_roots):
                    return path
        return None

    def _add_video_element(self, slide: Any, el: Dict[str, Any], x: Inches, y: Inches, w: Inches, h: Inches):
        content = el.get("content", "")
        video_path = self._resolve_local_video_path(content)
        if video_path:
            try:
                slide.shapes.add_movie(
                    str(video_path),
                    x,
                    y,
                    width=w,
                    height=h,
                    mime_type="video/mp4"
                )
                return
            except Exception as e:
                print(f"Error embedding video file in PPTX: {e}")

        # Fall back to visual placeholder if file resolution or embedding fails
        self._add_placeholder_element(slide, el, x, y, w, h)
