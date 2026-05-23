import json
import tempfile
from io import BytesIO
from pathlib import Path
from unittest import mock

import fitz  # type: ignore
from django.contrib.auth import get_user_model
from django.core.files.uploadedfile import SimpleUploadedFile
from django.test import Client, TestCase, override_settings
from PIL import Image
from pptx import Presentation as PptxPresentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from bridge.pdf_bridge import PDF2PPTxBridge
from bridge.processors import LocalVisionPDFProcessor
from bridge.text_utils import OutlineEntry
from bridge.pptx_exporter import PPTXExporter
from studio.models import Asset, Presentation, PresentationRevision

BACKEND_DIR = Path(__file__).resolve().parents[1]
PROJECT_ROOT = BACKEND_DIR.parent
FRONTEND_DIR = PROJECT_ROOT / "frontend"


class PresentationApiRegressionTests(TestCase):
    def setUp(self):
        self.client = Client()
        self.user = get_user_model().objects.create_user(username="tester", password="strong-password-123")
        self.client.force_login(self.user)

    def test_pptx_export_empty_deck_returns_pptx_response(self):
        response = self.client.post(
            "/api/presentations/export/pptx/",
            data=json.dumps({"state": {"pageSetup": "standard-4-3", "slides": []}, "filename": "test deck.pptx"}),
            content_type="application/json",
        )

        self.assertEqual(response.status_code, 200)
        self.assertEqual(
            response["Content-Type"],
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
        self.assertIn('filename="test_deck.pptx"', response["Content-Disposition"])

    def test_presentation_create_malformed_json_returns_400(self):
        response = self.client.post(
            "/api/presentations/",
            data="{bad json",
            content_type="application/json",
        )

        self.assertEqual(response.status_code, 400)

    def test_presentation_update_increments_version_and_revision_atomically(self):
        presentation = Presentation.objects.create(
            owner=self.user,
            title="Draft",
            state_json={"slides": []},
            autosave_version=3,
        )

        response = self.client.patch(
            f"/api/presentations/{presentation.id}/",
            data=json.dumps({"title": "Updated", "state": {"slides": [{"id": "s1"}]}, "saveRevision": True}),
            content_type="application/json",
        )

        self.assertEqual(response.status_code, 200)
        presentation.refresh_from_db()
        self.assertEqual(presentation.autosave_version, 4)
        self.assertEqual(presentation.title, "Updated")
        revision = PresentationRevision.objects.get(presentation=presentation)
        self.assertEqual(revision.version, 4)
        self.assertEqual(revision.state_json, {"slides": [{"id": "s1"}]})

    @mock.patch("studio.views.build_task_provider")
    @mock.patch("studio.views._generate_structured_json")
    @override_settings(PPTMAKER_ENABLE_AI_CLEANUP=True)
    def test_slide_cleanup_uses_llm_and_sanitizes_updates(self, generate_json, build_provider):
        build_provider.return_value = object()
        generate_json.return_value = {
            "summary": "Balanced headline and body.",
            "elements": [
                {
                    "id": "title",
                    "x": 48,
                    "y": 42,
                    "width": 820,
                    "height": 90,
                    "content": "Sharper title",
                    "styles": {"fontSize": "48px", "position": "fixed"},
                },
                {"id": "missing", "x": 0, "y": 0, "width": 10, "height": 10},
            ],
        }

        response = self.client.post(
            "/api/slides/cleanup/",
            data=json.dumps(
                {
                    "pageSetup": {"width": 1024, "height": 768},
                    "theme": "editorial",
                    "slide": {
                        "id": "slide-1",
                        "notes": "Make this crisp.",
                        "elements": [
                            {
                                "id": "title",
                                "type": "text",
                                "x": 100,
                                "y": 100,
                                "width": "500px",
                                "height": "80px",
                                "content": "Old title",
                                "styles": {"fontSize": "40px"},
                            }
                        ],
                    },
                }
            ),
            content_type="application/json",
        )

        self.assertEqual(response.status_code, 200)
        payload = response.json()
        self.assertEqual(payload["summary"], "Balanced headline and body.")
        self.assertEqual(len(payload["elements"]), 1)
        self.assertEqual(payload["elements"][0]["id"], "title")
        self.assertEqual(payload["elements"][0]["content"], "Sharper title")
        self.assertEqual(payload["elements"][0]["styles"], {"fontSize": "48px"})
        build_provider.assert_called_once_with(task="creative", allow_remote=True)
        generate_json.assert_called_once()

    def test_slide_cleanup_uses_local_cleanup_by_default(self):
        response = self.client.post(
            "/api/slides/cleanup/",
            data=json.dumps(
                {
                    "pageSetup": {"width": 1024, "height": 768},
                    "slide": {
                        "elements": [
                            {
                                "id": "title",
                                "type": "text",
                                "x": 103,
                                "y": 101,
                                "width": "500px",
                                "height": "80px",
                                "content": "Old title",
                            }
                        ],
                    },
                }
            ),
            content_type="application/json",
        )

        self.assertEqual(response.status_code, 200)
        payload = response.json()
        self.assertEqual(payload["fallback"], "local")
        self.assertEqual(payload["summary"], "Local cleanup applied")
        self.assertEqual(payload["elements"][0]["x"], 100)
        self.assertEqual(payload["elements"][0]["y"], 100)

    @mock.patch("studio.views.logger.warning")
    @mock.patch("studio.views.build_task_provider")
    @mock.patch("studio.views._generate_structured_json")
    @override_settings(PPTMAKER_ENABLE_AI_CLEANUP=True)
    def test_slide_cleanup_failure_uses_local_fallback(self, generate_json, build_provider, logger_warning):
        build_provider.return_value = object()
        generate_json.side_effect = RuntimeError("All LLM providers failed: secret provider detail")

        response = self.client.post(
            "/api/slides/cleanup/",
            data=json.dumps(
                {
                    "pageSetup": {"width": 1024, "height": 768},
                    "slide": {
                        "elements": [
                            {
                                "id": "title",
                                "type": "text",
                                "x": 100,
                                "y": 100,
                                "width": "500px",
                                "height": "80px",
                                "content": "Old title",
                            }
                        ],
                    },
                }
            ),
            content_type="application/json",
        )

        self.assertEqual(response.status_code, 200)
        payload = response.json()
        self.assertEqual(payload["fallback"], "ai_failed_local")
        self.assertEqual(payload["summary"], "Local cleanup applied after AI returned invalid JSON")
        self.assertEqual(payload["elements"], [])
        self.assertNotIn("secret provider detail", response.content.decode("utf-8"))
        build_provider.assert_called_once_with(task="creative", allow_remote=True)
        logger_warning.assert_called_once()

    @mock.patch("studio.views.logger.exception")
    def test_pptx_export_failure_returns_generic_error(self, logger_exception):
        with mock.patch("studio.views.PPTXExporter.export", side_effect=RuntimeError("internal detail")):
            response = self.client.post(
                "/api/presentations/export/pptx/",
                data=json.dumps({"state": {"pageSetup": "standard-4-3", "slides": []}}),
                content_type="application/json",
            )

        self.assertEqual(response.status_code, 500)
        self.assertEqual(response.json(), {"error": "PPTX export failed"})
        logger_exception.assert_called_once_with("PPTX export failed")


class PPTXExporterRegressionTests(TestCase):
    def test_image_stream_uses_configurable_project_root(self):
        with tempfile.TemporaryDirectory() as tmpdir:
            project_root = Path(tmpdir)
            media_dir = project_root / "media"
            media_dir.mkdir()
            image_path = media_dir / "sample.png"
            Image.new("RGB", (1, 1), color=(255, 0, 0)).save(image_path)

            exporter = PPTXExporter({"slides": []}, project_root=project_root)
            stream = exporter._get_image_stream("/media/sample.png")

            self.assertIsNotNone(stream)
            self.assertGreater(len(stream.getvalue()), 0)

    def test_image_stream_rejects_paths_outside_safe_asset_roots(self):
        with tempfile.TemporaryDirectory() as tmpdir:
            project_root = Path(tmpdir) / "project"
            project_root.mkdir()
            outside_dir = Path(tmpdir) / "outside"
            outside_dir.mkdir()
            image_path = outside_dir / "sample.png"
            Image.new("RGB", (1, 1), color=(255, 0, 0)).save(image_path)

            exporter = PPTXExporter({"slides": []}, project_root=project_root)
            self.assertIsNone(exporter._get_image_stream(str(image_path)))

    def test_css_font_weight_keywords_are_supported(self):
        exporter = PPTXExporter({"slides": []})

        self.assertEqual(exporter._parse_font_weight("normal"), 400)
        self.assertEqual(exporter._parse_font_weight("bold"), 700)
        self.assertEqual(exporter._parse_font_weight("600"), 600)

    def test_image_crop_transform_maps_to_pptx_crop_fractions(self):
        with tempfile.TemporaryDirectory() as tmpdir:
            project_root = Path(tmpdir)
            media_dir = project_root / "media"
            media_dir.mkdir()
            image_path = media_dir / "wide.png"
            Image.new("RGB", (400, 200), color=(255, 0, 0)).save(image_path)

            state = {
                "slides": [
                    {
                        "elements": [
                            {
                                "type": "image",
                                "x": 0,
                                "y": 0,
                                "width": "400px",
                                "height": "200px",
                                "content": "/media/wide.png",
                                "cropTransform": {
                                    "leftPercent": -50,
                                    "topPercent": -25,
                                    "widthPercent": 200,
                                    "heightPercent": 150,
                                },
                                "styles": {},
                            }
                        ]
                    }
                ]
            }

            prs = PptxPresentation(PPTXExporter(state, project_root=project_root).export())
            picture = next(shape for shape in prs.slides[0].shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE)

        self.assertAlmostEqual(picture.crop_left, 0.25, places=4)
        self.assertAlmostEqual(picture.crop_right, 0.25, places=4)
        self.assertAlmostEqual(picture.crop_top, 1 / 6, places=4)
        self.assertAlmostEqual(picture.crop_bottom, 1 / 6, places=4)

    def test_native_pptx_exports_tables_connectors_notes_and_placeholders(self):
        state = {
            "pageSetup": "widescreen-16-9",
            "slides": [
                {
                    "notes": "Speaker note",
                    "elements": [
                        {
                            "id": "title",
                            "type": "text",
                            "x": 40,
                            "y": 30,
                            "width": "420px",
                            "height": "100px",
                            "content": "Title <strong>bold</strong><br><em>italic</em>",
                            "styles": {"fontSize": "28px", "color": "#123456", "zIndex": 1},
                        },
                        {
                            "id": "tbl",
                            "type": "table",
                            "x": 60,
                            "y": 160,
                            "width": "520px",
                            "height": "180px",
                            "styles": {"zIndex": 2},
                            "tableData": {
                                "rows": 2,
                                "cols": 2,
                                "headerRow": True,
                                "cells": [[{"text": "Metric"}, {"text": "Value"}], [{"text": "Accuracy"}, {"text": "98%"}]],
                                "colWidths": [180, 140],
                                "rowHeights": [44, 52],
                            },
                        },
                        {
                            "id": "line",
                            "type": "connector",
                            "x": 650,
                            "y": 120,
                            "width": "160px",
                            "height": "80px",
                            "styles": {"color": "#ff0000", "strokeWidth": 3, "zIndex": 3},
                        },
                        {
                            "id": "html",
                            "type": "html",
                            "x": 720,
                            "y": 250,
                            "width": "220px",
                            "height": "120px",
                            "content": "<h1>Widget</h1>",
                            "styles": {"zIndex": 4},
                        },
                    ],
                }
            ],
        }

        prs = PptxPresentation(PPTXExporter(state).export())
        slide = prs.slides[0]

        self.assertEqual(slide.notes_slide.notes_text_frame.text, "Speaker note")
        self.assertTrue(any(shape.has_table for shape in slide.shapes))
        self.assertTrue(any(shape.shape_type == MSO_SHAPE_TYPE.LINE for shape in slide.shapes))
        all_text = "\n".join(shape.text for shape in slide.shapes if getattr(shape, "has_text_frame", False))
        self.assertIn("Title bold", all_text)
        self.assertIn("italic", all_text)
        self.assertIn("HTML embed placeholder", all_text)

        table = next(shape.table for shape in slide.shapes if shape.has_table)
        self.assertEqual(table.cell(0, 0).text, "Metric")
        self.assertEqual(table.cell(1, 1).text, "98%")

    def test_native_pptx_applies_theme_background_and_rgba_panel_colors(self):
        state = {
            "presentationTheme": "horizon",
            "slides": [
                {
                    "elements": [
                        {
                            "id": "panel-text",
                            "type": "text",
                            "x": 24,
                            "y": 32,
                            "width": "420px",
                            "height": "120px",
                            "content": [{"text": "Point one", "level": 0}],
                            "styles": {
                                "fontSize": "22px",
                                "color": "#eef4ff",
                                "backgroundColor": "rgba(255,255,255,0.07)",
                                "border": "1px solid rgba(125,211,252,0.40)",
                            },
                        }
                    ]
                }
            ],
        }

        prs = PptxPresentation(PPTXExporter(state).export())
        slide = prs.slides[0]
        text_shape = next(shape for shape in slide.shapes if getattr(shape, "has_text_frame", False) and "Point one" in shape.text)

        self.assertEqual(str(slide.background.fill.fore_color.rgb), "0B2442")
        self.assertEqual(str(text_shape.fill.fore_color.rgb), "1C334F")
        self.assertIn('buChar char="•"', text_shape.element.xml)

    def test_native_pptx_keeps_transparent_shapes_unfilled(self):
        state = {
            "presentationTheme": "editorial",
            "slides": [
                {
                    "elements": [
                        {
                            "id": "clear-box",
                            "type": "shape",
                            "shapeType": "rectangle",
                            "x": 20,
                            "y": 20,
                            "width": "200px",
                            "height": "100px",
                            "styles": {"backgroundColor": "transparent", "border": "2px solid #3B82F6"},
                        }
                    ]
                }
            ],
        }

        prs = PptxPresentation(PPTXExporter(state).export())
        shape = prs.slides[0].shapes[0]

        self.assertIn("<a:noFill/>", shape.element.xml)
        self.assertIn('val="3B82F6"', shape.element.xml)

    def test_native_pptx_exports_theme_aware_master_slide_elements(self):
        state = {
            "presentationTheme": "horizon",
            "masterSlides": {
                "content": {
                    "footerText": "Shared footer",
                    "logoText": "Lab Deck",
                    "showSlideNumber": True,
                }
            },
            "slides": [
                {
                    "masterId": "content",
                    "elements": [
                        {
                            "id": "body",
                            "type": "text",
                            "x": 80,
                            "y": 120,
                            "width": "500px",
                            "height": "80px",
                            "content": "Body text",
                            "styles": {"fontSize": "24px", "color": "#eef4ff"},
                        }
                    ],
                }
            ],
        }

        prs = PptxPresentation(PPTXExporter(state).export())
        slide_text = "\n".join(shape.text for shape in prs.slides[0].shapes if getattr(shape, "has_text_frame", False))

        self.assertIn("Shared footer", slide_text)
        self.assertIn("Lab Deck", slide_text)
        self.assertIn("01", slide_text)

    def test_native_pptx_preserves_inline_text_font_styles(self):
        state = {
            "slides": [
                {
                    "elements": [
                        {
                            "id": "styled-text",
                            "type": "text",
                            "x": 20,
                            "y": 20,
                            "width": "640px",
                            "height": "120px",
                            "content": (
                                'Normal <span style="font-size: 34px; font-family: Georgia, serif; '
                                'color: #ff0000; font-weight: 700; font-style: italic;">Styled</span>'
                            ),
                            "styles": {
                                "fontSize": "18px",
                                "fontFamily": "Arial, sans-serif",
                                "color": "#111111",
                            },
                        }
                    ]
                }
            ]
        }

        prs = PptxPresentation(PPTXExporter(state).export())
        text_shape = next(shape for shape in prs.slides[0].shapes if getattr(shape, "has_text_frame", False) and "Styled" in shape.text)
        runs = [run for paragraph in text_shape.text_frame.paragraphs for run in paragraph.runs]
        styled_run = next(run for run in runs if run.text == "Styled")

        self.assertEqual(styled_run.font.size.pt, 34)
        self.assertEqual(styled_run.font.name, "Georgia")
        self.assertEqual(str(styled_run.font.color.rgb), "FF0000")
        self.assertTrue(styled_run.font.bold)
        self.assertTrue(styled_run.font.italic)

    def test_native_pptx_preserves_block_text_font_styles(self):
        state = {
            "slides": [
                {
                    "elements": [
                        {
                            "id": "styled-block-text",
                            "type": "text",
                            "x": 20,
                            "y": 20,
                            "width": "640px",
                            "height": "120px",
                            "content": (
                                '<p style="font-size: 30px; font-family: Georgia, serif; '
                                'color: #336699; font-style: italic;">Block styled</p>'
                            ),
                            "styles": {
                                "fontSize": "18px",
                                "fontFamily": "Arial, sans-serif",
                                "color": "#111111",
                            },
                        }
                    ]
                }
            ]
        }

        prs = PptxPresentation(PPTXExporter(state).export())
        text_shape = next(shape for shape in prs.slides[0].shapes if getattr(shape, "has_text_frame", False) and "Block styled" in shape.text)
        run = next(run for paragraph in text_shape.text_frame.paragraphs for run in paragraph.runs if run.text == "Block styled")

        self.assertEqual(run.font.size.pt, 30)
        self.assertEqual(run.font.name, "Georgia")
        self.assertEqual(str(run.font.color.rgb), "336699")
        self.assertTrue(run.font.italic)

    def test_native_pptx_preserves_table_cell_font_styles(self):
        state = {
            "slides": [
                {
                    "elements": [
                        {
                            "id": "styled-table",
                            "type": "table",
                            "x": 20,
                            "y": 20,
                            "width": "420px",
                            "height": "120px",
                            "tableData": {
                                "rows": 1,
                                "cols": 1,
                                "headerRow": False,
                                "fontSize": "12px",
                                "fontFamily": "Arial, sans-serif",
                                "textColor": "#111111",
                                "cells": [
                                    [
                                        {
                                            "text": "Cell",
                                            "styles": {
                                                "fontSize": "28px",
                                                "fontFamily": "Georgia, serif",
                                                "fontWeight": "700",
                                                "fontStyle": "italic",
                                                "color": "#00aa00",
                                            },
                                        }
                                    ]
                                ],
                            },
                        }
                    ]
                }
            ]
        }

        prs = PptxPresentation(PPTXExporter(state).export())
        table = next(shape.table for shape in prs.slides[0].shapes if shape.has_table)
        run = table.cell(0, 0).text_frame.paragraphs[0].runs[0]

        self.assertEqual(run.font.size.pt, 28)
        self.assertEqual(run.font.name, "Georgia")
        self.assertEqual(str(run.font.color.rgb), "00AA00")
        self.assertTrue(run.font.bold)
        self.assertTrue(run.font.italic)


class FrontendExportRegressionTests(TestCase):
    def setUp(self):
        self.export_js = (FRONTEND_DIR / "js" / "export" / "export.js").read_text(encoding="utf-8")
        self.editor_export_block = "\n".join(self.export_js.splitlines()[:350])

    def test_pdf_and_png_exports_use_page_setup_dimensions_not_hardcoded_4_3(self):
        self.assertIn("getPresentationPageSetupConfig", self.editor_export_block)
        self.assertNotIn("format: [1024, 768]", self.editor_export_block)
        self.assertNotIn("width: 1024,", self.editor_export_block)
        self.assertNotIn("height: 768", self.editor_export_block)
        self.assertNotIn("pdf.addPage([1024, 768]", self.editor_export_block)
        self.assertNotIn('pdf.addImage(imgData, "JPEG", 0, 0, 1024, 768)', self.editor_export_block)

    def test_pdf_and_png_capture_active_slide_with_finally_restore(self):
        self.assertIn("getActiveExportSlideElement", self.export_js)
        self.assertIn("try {", self.export_js)
        self.assertIn("finally {", self.export_js)
        self.assertNotIn('document.getElementById("slides-container")', self.export_js)

    def test_zip_viewer_preserves_table_dimensions_and_wraps_html_embeds(self):
        self.assertIn("rawRowHeights", self.export_js)
        self.assertIn("rawColWidths", self.export_js)
        self.assertIn("document.createElement('colgroup')", self.export_js)
        self.assertIn("buildViewerHtmlEmbedSrcdoc", self.export_js)
        self.assertIn("iframe.setAttribute('sandbox'", self.export_js)

    def test_cropped_images_render_consistently_in_previews_editor_and_zip_viewer(self):
        render_js = (FRONTEND_DIR / "js" / "editor" / "render.js").read_text(encoding="utf-8")
        crop_js = (FRONTEND_DIR / "js" / "editor" / "crop.js").read_text(encoding="utf-8")

        self.assertIn("function _createImageContentNode", render_js)
        self.assertIn("_createImageContentNode(elData, { interactive: false })", render_js)
        self.assertIn("_createImageContentNode(elData, { interactive: true })", render_js)
        self.assertIn("img.style.setProperty(\"margin\", \"0\", \"important\")", render_js)
        self.assertIn("img.style.setProperty(\"margin\", \"0\", \"important\")", crop_js)
        self.assertIn("function normalizeImageCropTransform(crop)", self.export_js)
        self.assertIn("const crop = normalizeImageCropTransform(elData.cropTransform);", self.export_js)
        self.assertIn("margin:0!important", self.export_js)
        self.assertIn("max-height:none; object-fit:fill", self.export_js)

    def test_html_embeds_are_sandboxed_without_same_origin(self):
        html_embed_js = (FRONTEND_DIR / "js" / "embeds" / "htmlEmbed.js").read_text(encoding="utf-8")
        render_js = (FRONTEND_DIR / "js" / "editor" / "render.js").read_text(encoding="utf-8")
        properties_js = (FRONTEND_DIR / "js" / "properties.js").read_text(encoding="utf-8")

        self.assertIn("function applyHtmlEmbedSandbox", html_embed_js)
        self.assertIn('HTML_EMBED_SANDBOX = "allow-scripts allow-forms allow-popups allow-downloads"', html_embed_js)
        self.assertNotIn("allow-same-origin", html_embed_js)
        self.assertIn("applyHtmlEmbedSandbox(frame)", render_js)
        self.assertIn("applyHtmlEmbedSandbox(iframe)", render_js)
        self.assertIn("applyHtmlEmbedSandbox(frame)", properties_js)
        self.assertIn("getViewerHtmlEmbedSandbox", self.export_js)
        self.assertIn("iframe.setAttribute('referrerpolicy', 'no-referrer')", self.export_js)
        self.assertNotIn("allow-same-origin", self.export_js)


class StaticHardeningSourceTests(TestCase):
    def test_backend_source_routes_are_not_publicly_served(self):
        urls_py = (BACKEND_DIR / "pptmaker_backend" / "urls.py").read_text(encoding="utf-8")
        self.assertNotIn('r"^bridge/', urls_py)
        self.assertNotIn('r"^extracted_figures/', urls_py)
        self.assertIn("if settings.DEBUG:", urls_py)

    def test_frontend_cleanup_todos_are_applied(self):
        index_html = (FRONTEND_DIR / "index.html").read_text(encoding="utf-8")
        main_js = (FRONTEND_DIR / "js" / "core" / "main.js").read_text(encoding="utf-8")
        state_js = (FRONTEND_DIR / "js" / "core" / "state.js").read_text(encoding="utf-8")

        self.assertEqual(index_html.count("katex@0.16.9/dist/katex.min.css"), 1)
        self.assertIn('for="present-chalk-color-chip"', index_html)
        self.assertNotIn("console.log(`Grouped", main_js)
        self.assertNotIn("console.log(`Ungrouped", main_js)
        self.assertIn("if (modal) {", state_js)

    def test_imported_state_and_text_rendering_are_sanitized(self):
        state_js = (FRONTEND_DIR / "js" / "core" / "state.js").read_text(encoding="utf-8")
        text_content_js = (FRONTEND_DIR / "js" / "text" / "textContent.js").read_text(encoding="utf-8")

        self.assertIn("function sanitizeTextHtml", state_js)
        self.assertIn("SAFE_TEXT_TAGS", state_js)
        self.assertIn("SAFE_ELEMENT_TYPES", state_js)
        self.assertIn("sanitizeElementStyles", state_js)
        self.assertIn("sanitizeElementContent", state_js)
        self.assertIn("MAX_PRESENTATION_SLIDES", state_js)
        self.assertIn("normalizeStateIds();", state_js)
        self.assertIn("sanitizeTextHtml(elData.content", text_content_js)
        self.assertIn("sanitizeTextHtml(html)", text_content_js)


class AssetUploadHardeningTests(TestCase):
    def setUp(self):
        self.client = Client()
        self.user = get_user_model().objects.create_user(username="asset-user", password="strong-password-123")

    def _png_upload(self, name="sample.png"):
        image_bytes = BytesIO()
        Image.new("RGB", (2, 2), color=(0, 128, 255)).save(image_bytes, format="PNG")
        return SimpleUploadedFile(name, image_bytes.getvalue(), content_type="image/png")

    def test_asset_upload_requires_authentication(self):
        response = self.client.post("/api/assets/upload/", {"file": self._png_upload()})

        self.assertEqual(response.status_code, 401)
        self.assertFalse(Asset.objects.exists())

    @override_settings(PPTMAKER_MAX_ASSET_UPLOAD_BYTES=4)
    def test_asset_upload_rejects_oversized_file(self):
        self.client.force_login(self.user)
        response = self.client.post("/api/assets/upload/", {"file": self._png_upload()})

        self.assertEqual(response.status_code, 413)
        self.assertEqual(response.json(), {"error": "File is too large"})

    def test_asset_upload_rejects_invalid_image_content(self):
        self.client.force_login(self.user)
        upload = SimpleUploadedFile("sample.png", b"not an image", content_type="image/png")

        response = self.client.post("/api/assets/upload/", {"file": upload})

        self.assertEqual(response.status_code, 400)
        self.assertEqual(response.json(), {"error": "Uploaded image is not a valid image file"})

    def test_asset_upload_accepts_valid_image_for_authenticated_user(self):
        self.client.force_login(self.user)

        response = self.client.post("/api/assets/upload/", {"file": self._png_upload()})

        self.assertEqual(response.status_code, 201)
        self.assertEqual(response.json()["assetType"], "image")
        self.assertEqual(Asset.objects.get().owner, self.user)


class DocumentProcessingRegressionTests(TestCase):
    def test_embedded_pdf_visual_fallback_extracts_image_with_caption(self):
        with tempfile.TemporaryDirectory() as tmpdir:
            tmp_path = Path(tmpdir)
            image_path = tmp_path / "figure.png"
            Image.new("RGB", (180, 90), color=(32, 96, 160)).save(image_path)

            pdf_path = tmp_path / "source.pdf"
            doc = fitz.open()
            page = doc.new_page(width=420, height=320)
            page.insert_image(fitz.Rect(72, 60, 252, 150), filename=str(image_path))
            page.insert_textbox(
                fitz.Rect(72, 160, 360, 220),
                "Figure 1. Test caption describes the embedded result.",
                fontsize=11,
            )
            doc.save(pdf_path)
            doc.close()

            processor = LocalVisionPDFProcessor(str(pdf_path))
            processor.figures_dir = str(tmp_path / "extracted_figures")

            visuals = processor._extract_visuals_embedded()
            processor._sync_visual_context(visuals)

            self.assertEqual(len(visuals), 1)
            self.assertTrue(Path(visuals[0]["path"]).exists())
            self.assertIn("Figure 1", visuals[0]["caption"])
            self.assertIn("Test caption", processor.mineru_context["captions"][0])

    def test_storyboard_json_failure_uses_deterministic_fallback_with_debug(self):
        class BrokenLLM:
            trace = [{"provider": "test-provider", "error": "bad json"}]
            last_provider = "test-provider"

            def generate(self, *args, **kwargs):
                return '{"title": "Broken", "slides": [{"title": "unterminated}'

        bridge = PDF2PPTxBridge(BrokenLLM())
        events = []
        paper_brief = {
            "central_thesis": "String methods estimate transition pathways from swarms of trajectories.",
            "problem": "Free-energy calculations need tractable transition pathway estimates.",
            "method_mechanism": ["Initialize images along a trial path", "Evolve swarms from each image"],
            "key_findings": ["The method updates images toward average swarm drift"],
            "takeaway": "The workflow turns trajectory ensembles into interpretable pathways.",
            "evidence_items": [
                {
                    "id": "E1",
                    "claim": "Trial paths are represented by images.",
                    "support": "The source describes initializing images along the path.",
                    "section": "Method",
                    "figure_ids": ["F1"],
                },
                {
                    "id": "E2",
                    "claim": "Swarms estimate local drift.",
                    "support": "Short trajectories are launched from each image.",
                    "section": "Method",
                },
            ],
        }

        storyboard = bridge._plan_storyboard(
            BrokenLLM(),
            {"title": "String Method Tutorial"},
            [{"id": "F1", "caption": "String method schematic"}],
            lambda event, message, data=None: events.append((event, message, data or {})),
            document_outline=[
                OutlineEntry("Introduction", "Introduction", 1, "Transition pathways motivate the method.", 0, 10)
            ],
            coarse_context=[{"heading": "Introduction", "level": 1, "preview": "Transition pathways motivate the method."}],
            paper_brief=paper_brief,
        )

        self.assertGreaterEqual(len(storyboard["slides"]), 3)
        self.assertTrue(storyboard["_debug"]["storyboard"]["fallback"])
        self.assertIn("Failed to parse scientific storyboarding JSON", storyboard["_debug"]["storyboard"]["reason"])
        self.assertTrue(any(item[2].get("fallback") for item in events))

    def test_local_slide_generation_defaults_to_fast_source_writer(self):
        class LocalProvider:
            def provider_label(self):
                return "ollama:qwen2.5:7b"

        class LocalLLM:
            providers = [LocalProvider()]

            def generate(self, *args, **kwargs):
                raise AssertionError("fast local slide writing should not call the LLM")

        bridge = PDF2PPTxBridge(LocalLLM())
        storyboard = {
            "title": "String Method Tutorial",
            "slides": [
                {
                    "section": "Method",
                    "title": "Swarms estimate local drift",
                    "goal": "Short trajectories estimate drift around each image.",
                    "evidence_ids": ["E1"],
                    "layout_hint": "mechanism",
                    "context_keywords": ["swarms", "trajectories", "drift"],
                }
            ],
        }
        paper_brief = {
            "evidence_items": [
                {
                    "id": "E1",
                    "claim": "Short trajectories are launched from each image.",
                    "support": "The swarm average estimates local drift for path updates.",
                    "section": "Method",
                }
            ]
        }
        full_text = (
            "Methods. Short trajectories are launched from each image. "
            "The swarm average estimates local drift for path updates. "
            "The images move toward the average endpoint of the local trajectories."
        )

        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = Path(tmpdir) / "presentation.json"
            result = bridge._generate_slides(
                storyboard,
                [],
                full_text,
                {},
                LocalLLM(),
                output_path,
                None,
                0,
                0,
                lambda *args, **kwargs: None,
                {},
                {},
                paper_brief=paper_brief,
                llm_trace={"slide_writing": LocalLLM()},
            )

        self.assertEqual(result["generation"]["slide_writing_mode"], "fast")
        content_slides = [slide for slide in result["slides"] if slide.get("type") == "content"]
        self.assertEqual(len(content_slides), 1)
        self.assertEqual(content_slides[0]["review"]["status"], "fast_source_writer")
