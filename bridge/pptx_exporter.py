import io
import base64
import re
from html import unescape
from html.parser import HTMLParser
from pathlib import Path
from typing import Any, Dict, Optional
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN

class DesignSystem:
    # Default design constants inspired by the reference script
    DEFAULT_FONT = "Arial"
    ACCENT_COLOR = RGBColor(0, 51, 102)  # Dark Blue
    TEXT_COLOR = RGBColor(44, 62, 80)
    BG_COLOR = RGBColor(255, 255, 255)


class HtmlTextRunParser(HTMLParser):
    BLOCK_TAGS = {"div", "p", "li", "section", "article", "blockquote", "h1", "h2", "h3", "h4", "h5", "h6"}

    def __init__(self):
        super().__init__(convert_charrefs=True)
        self.runs: list[dict[str, Any]] = []
        self._style_stack: list[dict[str, bool]] = [{"bold": False, "italic": False, "underline": False}]

    @property
    def current_style(self) -> dict[str, bool]:
        return self._style_stack[-1]

    def _push(self, **updates):
        self._style_stack.append({**self.current_style, **updates})

    def _pop(self):
        if len(self._style_stack) > 1:
            self._style_stack.pop()

    def _newline(self):
        if self.runs and self.runs[-1].get("text") == "\n":
            return
        self.runs.append({"text": "\n", **self.current_style})

    def handle_starttag(self, tag, attrs):
        tag = tag.lower()
        if tag in {"strong", "b"}:
            self._push(bold=True)
        elif tag in {"em", "i"}:
            self._push(italic=True)
        elif tag == "u":
            self._push(underline=True)
        elif tag == "br":
            self._newline()
        elif tag in self.BLOCK_TAGS and self.runs:
            self._newline()

    def handle_endtag(self, tag):
        tag = tag.lower()
        if tag in {"strong", "b", "em", "i", "u"}:
            self._pop()
        elif tag in self.BLOCK_TAGS:
            self._newline()

    def handle_data(self, data):
        if data:
            self.runs.append({"text": unescape(data), **self.current_style})


def html_to_text_runs(content: Any) -> list[dict[str, Any]]:
    parser = HtmlTextRunParser()
    parser.feed(str(content or ""))
    runs = []
    for run in parser.runs:
        text = run.get("text", "")
        if text == "\n":
            if runs and runs[-1].get("text") == "\n":
                continue
            runs.append(run)
        elif text:
            runs.append(run)
    while runs and runs[-1].get("text") == "\n":
        runs.pop()
    return runs or [{"text": ""}]

class PPTXExporter:
    def __init__(self, state: Dict[str, Any], project_root: Optional[Path | str] = None):
        self.state = state
        self.prs = Presentation()
        self.project_root = self._resolve_project_root(project_root)
        
        # Mapping from SlideForge page setup to PPTX dimensions
        setup = state.get("pageSetup", "standard-4-3")
        if setup == "widescreen-16-9":
            self.prs.slide_width = Inches(13.333)
            self.prs.slide_height = Inches(7.5)
            self.base_w, self.base_h = 1280, 720
        elif setup == "widescreen-16-10":
            self.prs.slide_width = Inches(12.8)
            self.prs.slide_height = Inches(8.0)
            self.base_w, self.base_h = 1280, 800
        else: # standard-4-3
            self.prs.slide_width = Inches(10.666)
            self.prs.slide_height = Inches(8.0)
            self.base_w, self.base_h = 1024, 768

    def _resolve_project_root(self, project_root: Optional[Path | str] = None) -> Path:
        if project_root:
            return Path(project_root).expanduser().resolve()
        try:
            from django.conf import settings

            base_dir = getattr(settings, "BASE_DIR", None)
            if base_dir:
                return Path(base_dir).expanduser().resolve()
        except Exception:
            pass
        return Path(__file__).resolve().parent.parent

    def _safe_asset_roots(self) -> list[Path]:
        roots = [
            self.project_root / "media",
            self.project_root / "static",
            self.project_root / "staticfiles",
            self.project_root / "assets",
        ]
        try:
            from django.conf import settings

            for setting_name in ("MEDIA_ROOT", "STATIC_ROOT"):
                configured = getattr(settings, setting_name, None)
                if configured:
                    roots.append(Path(configured))
            for static_dir in getattr(settings, "STATICFILES_DIRS", []) or []:
                roots.append(Path(static_dir))
        except Exception:
            pass

        resolved = []
        seen = set()
        for root in roots:
            try:
                path = Path(root).expanduser().resolve()
            except Exception:
                continue
            if path in seen:
                continue
            seen.add(path)
            resolved.append(path)
        return resolved

    def _resolve_safe_local_image_path(self, content: str) -> Optional[Path]:
        normalized = str(content or "").strip()
        if not normalized or "://" in normalized:
            return None

        candidates: list[Path] = []
        if normalized.startswith("/media/"):
            candidates.append(self.project_root / normalized.lstrip("/"))
        elif normalized.startswith("/static/"):
            candidates.append(self.project_root / normalized.lstrip("/"))
            candidates.append(self.project_root / "staticfiles" / normalized.removeprefix("/static/"))
        elif normalized.startswith("/assets/"):
            candidates.append(self.project_root / normalized.lstrip("/"))
        elif normalized.startswith("/"):
            candidates.append(Path(normalized))
        else:
            candidates.append(self.project_root / normalized)

        safe_roots = self._safe_asset_roots()
        for candidate in candidates:
            try:
                path = candidate.expanduser().resolve()
            except Exception:
                continue
            if not path.exists() or not path.is_file():
                continue
            if any(path == root or path.is_relative_to(root) for root in safe_roots):
                return path
        return None

    def _parse_font_weight(self, value: Any) -> int:
        if isinstance(value, (int, float)):
            return int(value)
        normalized = str(value or "400").strip().lower()
        if normalized == "normal":
            return 400
        if normalized == "bold":
            return 700
        match = re.search(r"\d+", normalized)
        if match:
            return int(match.group(0))
        return 400

    def _px_to_inches_w(self, px: float) -> Inches:
        # Scale pixels relative to the logical canvas width
        return Inches((px / self.base_w) * (self.prs.slide_width / Inches(1)))

    def _px_to_inches_h(self, px: float) -> Inches:
        # Scale pixels relative to the logical canvas height
        return Inches((px / self.base_h) * (self.prs.slide_height / Inches(1)))

    def _parse_color(self, color_str: Optional[str], default: RGBColor = DesignSystem.TEXT_COLOR) -> RGBColor:
        if not color_str:
            return default
        try:
            if color_str.startswith("#"):
                h = color_str.lstrip("#")
                if len(h) == 3:
                    h = "".join([c*2 for c in h])
                return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
            elif color_str.startswith("rgb"):
                # Simple rgb(r, g, b) parser
                matches = re.findall(r"\d+", color_str)
                if len(matches) >= 3:
                    return RGBColor(int(matches[0]), int(matches[1]), int(matches[2]))
        except Exception:
            pass
        return default

    def _parse_px(self, value: Any, fallback: float = 100) -> float:
        match = re.search(r"-?\d+(?:\.\d+)?", str(value or ""))
        if not match:
            return fallback
        try:
            return float(match.group(0))
        except ValueError:
            return fallback

    def _set_text_common(self, run: Any, styles: dict[str, Any], *, extra: Optional[dict[str, bool]] = None):
        extra = extra or {}
        font_size = int(self._parse_px(styles.get("fontSize", "18"), 18))
        font_color = self._parse_color(styles.get("color"))
        font_name = str(styles.get("fontFamily", DesignSystem.DEFAULT_FONT)).split(",")[0].replace('"', "").strip()
        run.font.size = Pt(font_size)
        run.font.color.rgb = font_color
        run.font.name = font_name
        run.font.bold = self._parse_font_weight(styles.get("fontWeight")) >= 600 or extra.get("bold", False)
        run.font.italic = styles.get("fontStyle") == "italic" or extra.get("italic", False)
        run.font.underline = bool(extra.get("underline", False))

    def _paragraph_alignment(self, styles: dict[str, Any]):
        align = str(styles.get("textAlign", "left")).lower()
        if align == "center":
            return PP_ALIGN.CENTER
        if align == "right":
            return PP_ALIGN.RIGHT
        if align == "justify":
            return PP_ALIGN.JUSTIFY
        return PP_ALIGN.LEFT

    def _get_image_stream(self, content: str) -> Optional[io.BytesIO]:
        if not content:
            return None
            
        if content.startswith("data:"):
            try:
                base64_data = content.split(",")[1]
                return io.BytesIO(base64.b64decode(base64_data))
            except Exception:
                return None
        
        try:
            path = self._resolve_safe_local_image_path(content)
            if path:
                with open(path, "rb") as f:
                    return io.BytesIO(f.read())
        except Exception as e:
            print(f"Error resolving image path {content}: {e}")
            
        return None

    def export(self) -> io.BytesIO:
        slides = self.state.get("slides", [])
        for slide_data in slides:
            self._add_slide(slide_data)
        
        output = io.BytesIO()
        self.prs.save(output)
        output.seek(0)
        return output

    def _add_slide(self, slide_data: Dict[str, Any]):
        # Use a blank layout (index 6 is usually blank)
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 1. Background
        bg = slide_data.get("background")
        if bg and bg.get("content"):
            if bg.get("type") == "color":
                fill = slide.background.fill
                fill.solid()
                fill.fore_color.rgb = self._parse_color(bg["content"], DesignSystem.BG_COLOR)
            elif bg.get("type") == "image":
                img_stream = self._get_image_stream(bg["content"])
                if img_stream:
                    slide.shapes.add_picture(img_stream, 0, 0, width=self.prs.slide_width, height=self.prs.slide_height)
        elif slide_data.get("backgroundColor"):
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = self._parse_color(slide_data.get("backgroundColor"), DesignSystem.BG_COLOR)
        
        if slide_data.get("notes"):
            slide.notes_slide.notes_text_frame.text = str(slide_data.get("notes") or "")

        # 2. Elements
        elements = slide_data.get("elements", [])
        # Sort by zIndex
        sorted_elements = sorted(elements, key=lambda e: e.get("styles", {}).get("zIndex", 0))
        
        for el in sorted_elements:
            self._add_element(slide, el)

    def _add_element(self, slide: Any, el: Dict[str, Any]):
        try:
            x = self._px_to_inches_w(float(el.get("x", 0)))
            y = self._px_to_inches_h(float(el.get("y", 0)))
            
            # Width and height can be "auto" or strings like "400px"
            w_raw = el.get("width", "100")
            h_raw = el.get("height", "100")
            
            w = self._px_to_inches_w(float(re.sub(r"[^\d.]", "", str(w_raw)) or 100))
            h = self._px_to_inches_h(float(re.sub(r"[^\d.]", "", str(h_raw)) or 100))

            el_type = el.get("type")
            if el_type == "text":
                self._add_text_element(slide, el, x, y, w, h)
            elif el_type == "image":
                self._add_image_element(slide, el, x, y, w, h)
            elif el_type == "shape":
                self._add_shape_element(slide, el, x, y, w, h)
            elif el_type == "table":
                self._add_table_element(slide, el, x, y, w, h)
            elif el_type == "connector":
                self._add_connector_element(slide, el, x, y, w, h)
            elif el_type in {"video", "html", "pdf", "molecule", "chart", "equation", "latex"}:
                self._add_placeholder_element(slide, el, x, y, w, h)
        except Exception as e:
            print(f"Error adding element {el.get('id')}: {e}")

    def _add_text_element(self, slide: Any, el: Dict[str, Any], x: Inches, y: Inches, w: Inches, h: Inches):
        styles = el.get("styles", {})
        content = el.get("content", "")
        
        # Create textbox
        shape = slide.shapes.add_textbox(x, y, w, h)
        tf = shape.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        
        # Handle structured bullet points or raw HTML
        tf.clear()
        if isinstance(content, list):
            for idx, item in enumerate(content):
                p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                p.level = max(0, min(8, int(item.get("level", 0) or 0)))
                p.alignment = self._paragraph_alignment(styles)
                raw = item.get("html") if isinstance(item, dict) and item.get("html") is not None else item.get("text", "")
                self._append_runs_to_paragraph(p, html_to_text_runs(raw), styles)
        else:
            p = tf.paragraphs[0]
            p.alignment = self._paragraph_alignment(styles)
            self._append_runs_to_paragraph(p, html_to_text_runs(content), styles)

    def _append_runs_to_paragraph(self, paragraph: Any, runs: list[dict[str, Any]], styles: dict[str, Any]):
        current = paragraph
        for run_data in runs:
            text = run_data.get("text", "")
            parts = text.split("\n")
            for idx, part in enumerate(parts):
                if idx > 0:
                    current = paragraph._parent.add_paragraph()
                    current.alignment = paragraph.alignment
                if not part:
                    continue
                run = current.add_run()
                run.text = part
                self._set_text_common(run, styles, extra=run_data)

    def _add_image_element(self, slide: Any, el: Dict[str, Any], x: Inches, y: Inches, w: Inches, h: Inches):
        content = el.get("content", "")
        img_stream = self._get_image_stream(content)
        if not img_stream:
            return

        pic = slide.shapes.add_picture(img_stream, x, y, width=w, height=h)
        self._apply_image_crop_transform(pic, el.get("cropTransform"))

    def _apply_image_crop_transform(self, pic: Any, crop: Optional[dict[str, Any]]):
        if not crop:
            return
        try:
            width_pct = max(100.0, float(crop.get("widthPercent", 100) or 100))
            height_pct = max(100.0, float(crop.get("heightPercent", 100) or 100))
            left_pct = min(0.0, max(100.0 - width_pct, float(crop.get("leftPercent", 0) or 0)))
            top_pct = min(0.0, max(100.0 - height_pct, float(crop.get("topPercent", 0) or 0)))
        except (TypeError, ValueError):
            return

        pic.crop_left = max(0.0, min(1.0, -left_pct / width_pct))
        pic.crop_right = max(0.0, min(1.0, (width_pct + left_pct - 100.0) / width_pct))
        pic.crop_top = max(0.0, min(1.0, -top_pct / height_pct))
        pic.crop_bottom = max(0.0, min(1.0, (height_pct + top_pct - 100.0) / height_pct))

    def _add_shape_element(self, slide: Any, el: Dict[str, Any], x: Inches, y: Inches, w: Inches, h: Inches):
        shape_type_str = el.get("shapeType", "rectangle")
        
        mapping = {
            "rectangle": MSO_SHAPE.RECTANGLE,
            "circle": MSO_SHAPE.OVAL,
            "ellipse": MSO_SHAPE.OVAL,
            "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
            "diamond": MSO_SHAPE.DIAMOND,
            "hexagon": MSO_SHAPE.HEXAGON,
            "parallelogram": MSO_SHAPE.PARALLELOGRAM,
            "arrow-right": MSO_SHAPE.RIGHT_ARROW,
            "arrow-left": MSO_SHAPE.LEFT_ARROW,
            "arrow-up": MSO_SHAPE.UP_ARROW,
            "arrow-down": MSO_SHAPE.DOWN_ARROW,
        }
        
        mso_type = mapping.get(shape_type_str, MSO_SHAPE.RECTANGLE)
        shape = slide.shapes.add_shape(mso_type, x, y, w, h)
        
        styles = el.get("styles", {})
        bg_color = self._parse_color(styles.get("backgroundColor"), DesignSystem.ACCENT_COLOR)
        
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        
        # Border
        border = styles.get("border")
        if border:
            # Simplified border parser (e.g., "2px solid #000")
            match = re.search(r"(\d+)px", border)
            if match:
                shape.line.width = Pt(int(match.group(1)))
            color_match = re.search(r"#[0-9a-fA-F]{3,6}", border)
            if color_match:
                shape.line.color.rgb = self._parse_color(color_match.group(0))
        else:
            shape.line.fill.background()

    def _add_table_element(self, slide: Any, el: Dict[str, Any], x: Inches, y: Inches, w: Inches, h: Inches):
        table_data = el.get("tableData") or {}
        rows = max(1, int(table_data.get("rows") or len(table_data.get("cells") or []) or 1))
        cols = max(1, int(table_data.get("cols") or max((len(row) for row in table_data.get("cells", []) if isinstance(row, list)), default=1)))
        rows = min(rows, 50)
        cols = min(cols, 20)
        table_shape = slide.shapes.add_table(rows, cols, x, y, w, h)
        table = table_shape.table
        cells = table_data.get("cells") if isinstance(table_data.get("cells"), list) else []
        header_row = table_data.get("headerRow", True)
        font_size = table_data.get("fontSize") or "14px"

        for col_idx, raw_width in enumerate(table_data.get("colWidths") or []):
            if col_idx < cols:
                table.columns[col_idx].width = self._px_to_inches_w(max(36, float(raw_width or 140)))
        for row_idx, raw_height in enumerate(table_data.get("rowHeights") or []):
            if row_idx < rows:
                table.rows[row_idx].height = self._px_to_inches_h(max(24, float(raw_height or 44)))

        for row_idx in range(rows):
            for col_idx in range(cols):
                cell = table.cell(row_idx, col_idx)
                raw_cell = cells[row_idx][col_idx] if row_idx < len(cells) and isinstance(cells[row_idx], list) and col_idx < len(cells[row_idx]) else {}
                text = raw_cell.get("text", "") if isinstance(raw_cell, dict) else str(raw_cell or "")
                cell.text = str(text)
                fill_color = (
                    table_data.get("headerFill", "#e2e8f0")
                    if row_idx == 0 and header_row
                    else table_data.get("altFill", "#f8fafc") if table_data.get("zebra") and row_idx % 2 == 0 else table_data.get("bodyFill", "#ffffff")
                )
                cell.fill.solid()
                cell.fill.fore_color.rgb = self._parse_color(fill_color, RGBColor(255, 255, 255))
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.alignment = self._paragraph_alignment({"textAlign": table_data.get("textAlign", "left")})
                    for run in paragraph.runs:
                        self._set_text_common(
                            run,
                            {
                                "fontSize": font_size,
                                "fontFamily": table_data.get("fontFamily", DesignSystem.DEFAULT_FONT),
                                "fontWeight": "700" if row_idx == 0 and header_row else table_data.get("fontWeight", "400"),
                                "color": table_data.get("headerTextColor" if row_idx == 0 and header_row else "textColor", "#172033"),
                            },
                        )
    def _add_connector_element(self, slide: Any, el: Dict[str, Any], x: Inches, y: Inches, w: Inches, h: Inches):
        styles = el.get("styles", {})
        connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x, y, x + w, y + h)
        connector.line.color.rgb = self._parse_color(styles.get("color") or styles.get("borderColor"), DesignSystem.ACCENT_COLOR)
        connector.line.width = Pt(float(styles.get("strokeWidth") or styles.get("borderWidth") or 2))

    def _add_placeholder_element(self, slide: Any, el: Dict[str, Any], x: Inches, y: Inches, w: Inches, h: Inches):
        label_by_type = {
            "video": "Video",
            "html": "HTML embed",
            "pdf": "PDF",
            "molecule": "3D molecule",
            "chart": "Chart",
            "equation": "Equation",
            "latex": "Equation",
        }
        label = label_by_type.get(el.get("type"), "Unsupported content")
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(241, 245, 249)
        shape.line.color.rgb = RGBColor(148, 163, 184)
        tf = shape.text_frame
        tf.text = f"{label} placeholder"
        paragraph = tf.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        if paragraph.runs:
            run = paragraph.runs[0]
        else:
            run = paragraph.add_run()
            run.text = f"{label} placeholder"
        run.font.size = Pt(14)
        run.font.color.rgb = RGBColor(71, 85, 105)
