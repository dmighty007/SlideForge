import json
import os
from pathlib import Path
from typing import Any, Dict, Iterable, List, Optional

import fitz  # type: ignore
import requests
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt


class DesignSystem:
    FONT = "Arial"
    HEADING = RGBColor(0, 0, 0)
    TEXT = RGBColor(60, 60, 60)
    ACCENT = RGBColor(0, 102, 204)
    BG_COLOR = RGBColor(255, 255, 255)


class LLMProvider:
    def generate(self, prompt: str, system_prompt: str = "", json_mode: bool = False) -> str:
        raise NotImplementedError


class OllamaProvider(LLMProvider):
    def __init__(self, model: str | None = None, models: Iterable[str] | None = None, base_url: str | None = None):
        configured = list(models or [])
        if model:
            configured.insert(0, model)
        self.models = [name for name in configured if name]
        self.base_url = (base_url or os.getenv("OLLAMA_BASE_URL", "http://127.0.0.1:11434")).rstrip("/")

    def generate(self, prompt: str, system_prompt: str = "", json_mode: bool = False) -> str:
        if not self.models:
            raise RuntimeError("No Ollama model configured")

        last_error: Exception | None = None
        for model_name in self.models:
            payload: Dict[str, Any] = {
                "model": model_name,
                "prompt": prompt,
                "system": system_prompt,
                "stream": False,
            }
            if json_mode:
                payload["format"] = "json"
            try:
                response = requests.post(f"{self.base_url}/api/generate", json=payload, timeout=300)
                response.raise_for_status()
                body = response.json()
                return str(body.get("response", ""))
            except Exception as exc:
                last_error = exc
        raise RuntimeError(f"Ollama request failed for all models: {last_error}") from last_error


class GroqProvider(LLMProvider):
    def __init__(self, api_key: str | None = None, model: str | None = None):
        self.api_key = api_key or os.getenv("GROQ_API_KEY", "")
        self.model = model or os.getenv("GROQ_MODEL", "llama-3.3-70b-versatile")

    def generate(self, prompt: str, system_prompt: str = "", json_mode: bool = False) -> str:
        if not self.api_key:
            raise RuntimeError("GROQ_API_KEY is not configured")
        payload: Dict[str, Any] = {
            "model": self.model,
            "messages": [
                {"role": "system", "content": system_prompt or "You are a precise assistant."},
                {"role": "user", "content": prompt},
            ],
        }
        if json_mode:
            payload["response_format"] = {"type": "json_object"}
        response = requests.post(
            "https://api.groq.com/openai/v1/chat/completions",
            headers={"Authorization": f"Bearer {self.api_key}", "Content-Type": "application/json"},
            json=payload,
            timeout=300,
        )
        response.raise_for_status()
        body = response.json()
        return body["choices"][0]["message"]["content"]


class GeminiProvider(LLMProvider):
    def __init__(self, api_key: str | None = None, model: str | None = None):
        self.api_key = api_key or os.getenv("GOOGLE_API_KEY", "")
        self.model = model or os.getenv("GEMINI_MODEL", "gemini-2.0-flash")

    def generate(self, prompt: str, system_prompt: str = "", json_mode: bool = False) -> str:
        if not self.api_key:
            raise RuntimeError("GOOGLE_API_KEY is not configured")
        parts = []
        if system_prompt:
            parts.append({"text": system_prompt})
        parts.append({"text": prompt})
        payload: Dict[str, Any] = {"contents": [{"parts": parts}]}
        if json_mode:
            payload["generationConfig"] = {"responseMimeType": "application/json"}
        response = requests.post(
            (
                "https://generativelanguage.googleapis.com/v1beta/models/"
                f"{self.model}:generateContent?key={self.api_key}"
            ),
            headers={"Content-Type": "application/json"},
            json=payload,
            timeout=300,
        )
        response.raise_for_status()
        body = response.json()
        candidates = body.get("candidates") or []
        if not candidates:
            raise RuntimeError(f"Gemini returned no candidates: {body}")
        parts_out = candidates[0].get("content", {}).get("parts", [])
        text = "".join(str(part.get("text", "")) for part in parts_out)
        if not text:
            raise RuntimeError(f"Gemini returned empty content: {body}")
        return text


class DeepSeekProvider(LLMProvider):
    def __init__(self, api_key: str | None = None, model: str | None = None):
        self.api_key = api_key or os.getenv("DEEPSEEK_API_KEY", "")
        self.model = model or os.getenv("DEEPSEEK_MODEL", "deepseek-chat")

    def generate(self, prompt: str, system_prompt: str = "", json_mode: bool = False) -> str:
        if not self.api_key:
            raise RuntimeError("DEEPSEEK_API_KEY is not configured")
        payload: Dict[str, Any] = {
            "model": self.model,
            "messages": [
                {"role": "system", "content": system_prompt or "You are a precise assistant."},
                {"role": "user", "content": prompt},
            ],
        }
        if json_mode:
            payload["response_format"] = {"type": "json_object"}
        response = requests.post(
            "https://api.deepseek.com/chat/completions",
            headers={"Authorization": f"Bearer {self.api_key}", "Content-Type": "application/json"},
            json=payload,
            timeout=300,
        )
        response.raise_for_status()
        body = response.json()
        return body["choices"][0]["message"]["content"]


class SmartFallbackProvider(LLMProvider):
    def __init__(self, providers: Iterable[LLMProvider]):
        self.providers = [provider for provider in providers if provider]

    def generate(self, prompt: str, system_prompt: str = "", json_mode: bool = False) -> str:
        if not self.providers:
            raise RuntimeError("No LLM providers configured")

        errors = []
        for provider in self.providers:
            try:
                return provider.generate(prompt, system_prompt, json_mode=json_mode)
            except Exception as exc:
                errors.append(f"{provider.__class__.__name__}: {exc}")
        raise RuntimeError("All LLM providers failed: " + " | ".join(errors))


class PDFProcessor:
    def __init__(self, filepath: str, llm_provider: LLMProvider | None = None):
        self.filepath = str(filepath)
        self.llm = llm_provider
        self.figures_dir = str(Path.cwd() / "extracted_figures")
        self.visual_catalog: List[Dict[str, Any]] = []
        self.mineru_context: Dict[str, Any] = {
            "backend": "marker",
            "visuals": [],
            "captions": [],
            "equations": [],
        }
        Path(self.figures_dir).mkdir(parents=True, exist_ok=True)

    def extract_visuals_hybrid(self) -> list[dict]:
        self.visual_catalog = []
        self.mineru_context["visuals"] = []
        self.mineru_context["captions"] = []
        return self.visual_catalog

    def extract_text(self, start: int = 0, end: int = -1) -> str:
        with fitz.open(self.filepath) as doc:
            last_page = len(doc) - 1
            if end < 0 or end > last_page:
                end = last_page
            start = max(0, min(start, last_page if last_page >= 0 else 0))
            if end < start:
                return ""
            chunks = [doc.load_page(page_idx).get_text("text") for page_idx in range(start, end + 1)]
        return "\n".join(chunk for chunk in chunks if chunk)


class PDF2PPTx:
    def __init__(self, llm_provider: LLMProvider | None = None):
        self.llm = llm_provider


class PPTXGenerator:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)

    def _apply_theme(self, slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = DesignSystem.BG_COLOR
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0,
            0,
            self.prs.slide_width,
            int(self.prs.slide_height * 0.012),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = DesignSystem.ACCENT
        bar.line.fill.background()

    def add_title(self, title: str, sub: str):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[0])
        self._apply_theme(slide)
        title_box = slide.shapes.title
        if title_box:
            title_box.text = title or "Untitled Presentation"
            paragraph = title_box.text_frame.paragraphs[0]
            paragraph.font.name = DesignSystem.FONT
            paragraph.font.bold = True
            paragraph.font.size = Pt(44)
            paragraph.font.color.rgb = DesignSystem.HEADING

        subtitle = slide.placeholders[1] if len(slide.placeholders) > 1 else None
        if subtitle:
            subtitle.text = sub or ""
            paragraph = subtitle.text_frame.paragraphs[0]
            paragraph.font.name = DesignSystem.FONT
            paragraph.font.color.rgb = DesignSystem.TEXT

    def add_section_slide(self, section_name: str):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[2])
        self._apply_theme(slide)
        title_box = slide.shapes.title
        if title_box:
            title_box.text = section_name or "Section"
            paragraph = title_box.text_frame.paragraphs[0]
            paragraph.font.name = DesignSystem.FONT
            paragraph.font.bold = True
            paragraph.font.color.rgb = DesignSystem.HEADING

    def add_content_slide(
        self,
        title: str,
        points: List[Dict[str, Any]],
        fig_path: str | None = None,
        fig_cap: str | None = None,
        fig_paths: List[str] | None = None,
    ):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        self._apply_theme(slide)
        title_box = slide.shapes.title
        if title_box:
            title_box.text = title or "Untitled"
            paragraph = title_box.text_frame.paragraphs[0]
            paragraph.font.name = DesignSystem.FONT
            paragraph.font.bold = True
            paragraph.font.size = Pt(28)
            paragraph.font.color.rgb = DesignSystem.HEADING

        body = slide.placeholders[1]
        text_frame = body.text_frame
        text_frame.clear()

        paths = fig_paths if fig_paths else ([fig_path] if fig_path else [])
        if paths:
            body.width = int(self.prs.slide_width * 0.48)
            available_height = int(self.prs.slide_height * 0.7)
            top_y = body.top
            left_x = int(self.prs.slide_width * 0.52)
            width = int(self.prs.slide_width * 0.45)

            valid_paths = [path for path in paths if path and os.path.exists(path)]
            fig_count = len(valid_paths)
            if fig_count > 0:
                each_height = available_height // fig_count
                for index, path in enumerate(valid_paths[:3]):
                    try:
                        with Image.open(path) as img:
                            aspect = img.height / img.width
                        height = min(int(width * aspect), each_height - Pt(10))
                        slide.shapes.add_picture(path, left_x, top_y + (index * each_height), width=width, height=height)
                    except Exception:
                        pass

        total_items = sum(
            len(item.get("content", [])) if isinstance(item.get("content", []), list) else 1 for item in points
        ) + len(points)
        heading_size, body_size = (Pt(16), Pt(12)) if total_items < 10 else (Pt(14), Pt(10))

        for item in points[:4]:
            paragraph = text_frame.add_paragraph()
            paragraph.text = item.get("heading", "")
            paragraph.font.name = DesignSystem.FONT
            paragraph.font.bold = True
            paragraph.font.size = heading_size
            paragraph.font.color.rgb = DesignSystem.ACCENT

            sub_bullets = item.get("content", [])
            if isinstance(sub_bullets, str):
                sub_bullets = [sub_bullets]
            for bullet in sub_bullets[:3]:
                subparagraph = text_frame.add_paragraph()
                subparagraph.text = str(bullet).strip()
                subparagraph.level = 1
                subparagraph.font.name = DesignSystem.FONT
                subparagraph.font.size = body_size
                subparagraph.font.color.rgb = DesignSystem.TEXT

    def save(self, path: str):
        self.prs.save(path)
